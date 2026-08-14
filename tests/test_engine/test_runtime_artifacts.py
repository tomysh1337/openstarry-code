from __future__ import annotations

import asyncio
import contextlib
import hashlib
import json
import os
import threading
from collections.abc import AsyncIterator
from io import BytesIO
from types import SimpleNamespace
from typing import Any

import pytest
from pptx import Presentation

import openstarry_code.engine.agent as agent_module
from openstarry_code.artifacts import ArtifactStore, artifact_payload
from openstarry_code.engine import Agent, AgentConfig, ToolCall, ToolResult
from openstarry_code.engine.artifact_delivery import (
    artifact_delivery_publish_target_key,
    auto_publish_omitted_workspace_artifacts,
)
from openstarry_code.engine.runtime import TurnRunner
from openstarry_code.engine.types import (
    ArtifactEvent,
    DoneEvent,
    ErrorEvent,
    RouterDecisionEvent,
    TextDeltaEvent,
    ToolResultEvent,
    ToolUseStartEvent,
)
from openstarry_code.gateway.config import AttachmentsConfig, GatewayConfig, SquillaRouterConfig
from openstarry_code.provider import (
    ContentBlockToolResult,
    Message,
    ModelInfo,
    ToolDefinition,
    ToolInputSchema,
)
from openstarry_code.provider import DoneEvent as ProviderDone
from openstarry_code.provider import ErrorEvent as ProviderError
from openstarry_code.provider import ReasoningDeltaEvent as ProviderReasoning
from openstarry_code.provider import TextDeltaEvent as ProviderText
from openstarry_code.provider import ToolUseEndEvent as ProviderToolUseEnd
from openstarry_code.provider import ToolUseStartEvent as ProviderToolUseStart
from openstarry_code.session.manager import SessionManager
from openstarry_code.session.storage import SessionStorage
from openstarry_code.tools.builtin import filesystem
from openstarry_code.tools.builtin import patch as patch_tools
from openstarry_code.tools.registry import ToolRegistry, ToolSpec
from openstarry_code.tools.types import (
    CallerKind,
    RetryableToolInputError,
    ToolContext,
    ToolError,
    current_tool_context,
)


class _ArtifactProvider:
    provider_name = "test"

    def __init__(self) -> None:
        self.calls = 0
        self.model = "test/model"

    def chat(self, messages: list[Message], tools=None, config=None) -> AsyncIterator[Any]:
        self.calls += 1
        return self._stream(self.calls)

    async def _stream(self, call_number: int) -> AsyncIterator[Any]:
        if call_number == 1:
            yield ProviderToolUseStart(tool_use_id="tool-1", tool_name="make_file")
            yield ProviderToolUseEnd(
                tool_use_id="tool-1",
                tool_name="make_file",
                arguments={},
            )
            yield ProviderDone(stop_reason="tool_use", input_tokens=1, output_tokens=1)
            return
        yield ProviderText(text="done")
        yield ProviderDone(stop_reason="stop", input_tokens=1, output_tokens=1)

    async def list_models(self) -> list[ModelInfo]:
        return []


class _BlockingAfterArtifactProvider(_ArtifactProvider):
    async def _stream(self, call_number: int) -> AsyncIterator[Any]:
        if call_number == 1:
            async for event in super()._stream(call_number):
                yield event
            return
        await asyncio.Event().wait()


class _PostPublishToolLoopProvider:
    provider_name = "test"

    def __init__(self) -> None:
        self.calls = 0
        self.model = "test/model"
        self.tools_seen: list[bool] = []

    def chat(self, messages: list[Message], tools=None, config=None) -> AsyncIterator[Any]:
        self.calls += 1
        self.tools_seen.append(bool(tools))
        return self._stream(self.calls)

    async def _stream(self, call_number: int) -> AsyncIterator[Any]:
        if call_number == 1:
            yield ProviderText(text="Preparing your presentation.")
            yield ProviderToolUseStart(
                tool_use_id="publish-1",
                tool_name="publish_artifact",
            )
            yield ProviderToolUseEnd(
                tool_use_id="publish-1",
                tool_name="publish_artifact",
                arguments={"path": "report.pptx"},
            )
            yield ProviderDone(stop_reason="tool_use", input_tokens=1, output_tokens=1)
            return
        yield ProviderToolUseStart(tool_use_id="qa-1", tool_name="qa_check")
        yield ProviderToolUseEnd(
            tool_use_id="qa-1",
            tool_name="qa_check",
            arguments={"path": "report.pptx"},
        )
        yield ProviderDone(stop_reason="tool_use", input_tokens=1, output_tokens=1)

    async def list_models(self) -> list[ModelInfo]:
        return []


class _GoalPostPublishLoopProvider:
    provider_name = "test"

    def __init__(
        self,
        *,
        plain_final: bool = False,
        post_publish_error_code: str | None = None,
        post_publish_tool_failure: bool = False,
        terminal_status: str = "complete",
        illegal_terminal_summary_tool: bool = False,
        terminal_summary_mode: str = "normal",
    ) -> None:
        self.plain_final = plain_final
        self.post_publish_error_code = post_publish_error_code
        self.post_publish_tool_failure = post_publish_tool_failure
        self.terminal_status = terminal_status
        self.illegal_terminal_summary_tool = illegal_terminal_summary_tool
        self.terminal_summary_mode = terminal_summary_mode
        self.calls = 0
        self.model = "test/model"
        self.tool_names_seen: list[list[str]] = []
        self.requests: list[list[Message]] = []

    def chat(self, messages: list[Message], tools=None, config=None) -> AsyncIterator[Any]:
        self.calls += 1
        self.requests.append(list(messages))
        self.tool_names_seen.append([tool.name for tool in tools or []])
        return self._stream(self.calls)

    async def _stream(self, call_number: int) -> AsyncIterator[Any]:
        if call_number == 1:
            yield ProviderToolUseStart(
                tool_use_id="publish-1",
                tool_name="publish_artifact",
            )
            yield ProviderToolUseEnd(
                tool_use_id="publish-1",
                tool_name="publish_artifact",
                arguments={"path": "report.html"},
            )
            yield ProviderDone(stop_reason="tool_use", input_tokens=1, output_tokens=1)
            return
        if call_number == 2 and self.post_publish_error_code is not None:
            yield ProviderError(
                message="Synthetic provider failure after durable artifact delivery.",
                code=self.post_publish_error_code,
            )
            return
        if call_number == 2 and self.post_publish_tool_failure:
            yield ProviderToolUseStart(tool_use_id="qa-2", tool_name="qa_check")
            yield ProviderToolUseEnd(
                tool_use_id="qa-2",
                tool_name="qa_check",
                arguments={"path": "report.html"},
            )
            yield ProviderDone(stop_reason="tool_use", input_tokens=1, output_tokens=1)
            return
        if call_number == 3 and self.post_publish_tool_failure:
            yield ProviderText(text="The published report remains available for review.")
            yield ProviderDone(stop_reason="stop", input_tokens=1, output_tokens=1)
            return
        if self.plain_final:
            yield ProviderText(text="The report is ready.")
            yield ProviderDone(stop_reason="stop", input_tokens=1, output_tokens=1)
            return
        if call_number == 2:
            yield ProviderToolUseStart(
                tool_use_id="progress-2",
                tool_name="update_goal_progress",
            )
            yield ProviderToolUseEnd(
                tool_use_id="progress-2",
                tool_name="update_goal_progress",
                arguments={
                    "steps": [
                        {
                            "step": "Publish the verified artifact",
                            "status": "completed",
                        }
                    ]
                },
            )
            yield ProviderDone(stop_reason="tool_use", input_tokens=1, output_tokens=1)
            return
        if call_number == 3:
            arguments = {"status": self.terminal_status}
            if self.terminal_status == "blocked":
                arguments["reason"] = "Synthetic blocker"
            yield ProviderToolUseStart(
                tool_use_id="goal-3",
                tool_name="update_goal",
            )
            yield ProviderToolUseEnd(
                tool_use_id="goal-3",
                tool_name="update_goal",
                arguments=arguments,
            )
            yield ProviderDone(stop_reason="tool_use", input_tokens=1, output_tokens=1)
            return
        if self.illegal_terminal_summary_tool:
            yield ProviderToolUseStart(tool_use_id="qa-illegal", tool_name="qa_check")
            yield ProviderToolUseEnd(
                tool_use_id="qa-illegal",
                tool_name="qa_check",
                arguments={"path": "report.html"},
            )
            yield ProviderDone(stop_reason="tool_use", input_tokens=1, output_tokens=1)
            return
        if self.terminal_summary_mode == "empty":
            yield ProviderDone(stop_reason="stop", input_tokens=1, output_tokens=0)
            return
        if self.terminal_summary_mode == "reasoning_only":
            yield ProviderDone(
                stop_reason="stop",
                input_tokens=1,
                output_tokens=1,
                reasoning_tokens=1,
                reasoning_content="Synthetic internal reasoning.",
            )
            return
        if self.terminal_summary_mode == "stream_incomplete":
            yield ProviderText(text="Partial terminal summary")
            return
        if self.terminal_summary_mode == "length_capped":
            yield ProviderText(text="Partial terminal summary")
            yield ProviderDone(stop_reason="length", input_tokens=1, output_tokens=1)
            return
        yield ProviderText(text="The Goal is complete.")
        yield ProviderDone(stop_reason="stop", input_tokens=1, output_tokens=1)

    async def list_models(self) -> list[ModelInfo]:
        return []


class _GoalPublishAndYieldProvider:
    provider_name = "test"

    def __init__(self) -> None:
        self.calls = 0
        self.model = "test/model"
        self.tools_by_call: list[list[str]] = []

    def chat(self, messages: list[Message], tools=None, config=None) -> AsyncIterator[Any]:
        self.calls += 1
        self.tools_by_call.append(
            [str(getattr(tool, "name", "")) for tool in tools or []]
        )
        return self._stream(self.calls)

    async def _stream(self, call_number: int) -> AsyncIterator[Any]:
        if call_number > 1:
            yield ProviderText(text="The Goal is complete and the artifact is ready.")
            yield ProviderDone(stop_reason="stop", input_tokens=1, output_tokens=1)
            return
        yield ProviderToolUseStart(
            tool_use_id="publish-1",
            tool_name="publish_artifact",
        )
        yield ProviderToolUseEnd(
            tool_use_id="publish-1",
            tool_name="publish_artifact",
            arguments={"path": "report.html"},
        )
        yield ProviderToolUseStart(
            tool_use_id="goal-1",
            tool_name="update_goal",
        )
        yield ProviderToolUseEnd(
            tool_use_id="goal-1",
            tool_name="update_goal",
            arguments={"status": "complete"},
        )
        yield ProviderToolUseStart(
            tool_use_id="yield-1",
            tool_name="sessions_yield",
        )
        yield ProviderToolUseEnd(
            tool_use_id="yield-1",
            tool_name="sessions_yield",
            arguments={},
        )
        yield ProviderDone(stop_reason="tool_use", input_tokens=1, output_tokens=1)

    async def list_models(self) -> list[ModelInfo]:
        return []


class _GoalTerminalThenPublishProvider:
    provider_name = "test"

    def __init__(self, *, summary_mode: str = "normal") -> None:
        self.calls = 0
        self.model = "test/model"
        self.summary_mode = summary_mode
        self.tools_by_call: list[list[str]] = []

    def chat(self, messages: list[Message], tools=None, config=None) -> AsyncIterator[Any]:
        self.calls += 1
        self.tools_by_call.append(
            [str(getattr(tool, "name", "")) for tool in tools or []]
        )
        return self._stream(self.calls)

    async def _stream(self, call_number: int) -> AsyncIterator[Any]:
        if call_number == 1:
            yield ProviderToolUseStart(
                tool_use_id="goal-terminal-first",
                tool_name="update_goal",
            )
            yield ProviderToolUseEnd(
                tool_use_id="goal-terminal-first",
                tool_name="update_goal",
                arguments={"status": "complete"},
            )
            yield ProviderToolUseStart(
                tool_use_id="publish-after-terminal",
                tool_name="publish_artifact",
            )
            yield ProviderToolUseEnd(
                tool_use_id="publish-after-terminal",
                tool_name="publish_artifact",
                arguments={"path": "late.html"},
            )
            yield ProviderDone(stop_reason="tool_use", input_tokens=1, output_tokens=1)
            return
        if self.summary_mode == "reasoning_stream":
            yield ProviderReasoning(text="r" * 20)
            yield ProviderDone(
                stop_reason="stop",
                input_tokens=1,
                output_tokens=1,
                reasoning_tokens=20,
                reasoning_content="r" * 20,
            )
            return
        if self.summary_mode == "thinking_error":
            yield ProviderError(
                message="Synthetic thinking mode failure.",
                code="synthetic_thinking_failure",
            )
            return
        yield ProviderText(text="Final Goal summary.")
        yield ProviderDone(stop_reason="stop", input_tokens=1, output_tokens=1)

    async def list_models(self) -> list[ModelInfo]:
        return []


class _GoalYieldThenTerminalProvider:
    provider_name = "test"

    def __init__(self) -> None:
        self.calls = 0
        self.model = "test/model"

    def chat(self, messages: list[Message], tools=None, config=None) -> AsyncIterator[Any]:
        del messages, tools, config
        self.calls += 1
        return self._stream()

    async def _stream(self) -> AsyncIterator[Any]:
        yield ProviderToolUseStart(tool_use_id="yield-first", tool_name="sessions_yield")
        yield ProviderToolUseEnd(
            tool_use_id="yield-first",
            tool_name="sessions_yield",
            arguments={},
        )
        yield ProviderToolUseStart(tool_use_id="goal-after-yield", tool_name="update_goal")
        yield ProviderToolUseEnd(
            tool_use_id="goal-after-yield",
            tool_name="update_goal",
            arguments={"status": "complete"},
        )
        yield ProviderDone(stop_reason="tool_use", input_tokens=1, output_tokens=1)

    async def list_models(self) -> list[ModelInfo]:
        return []


class _SelectorClone:
    current_config = SimpleNamespace(model="test/model")

    def __init__(self, provider: _ArtifactProvider) -> None:
        self.provider = provider

    def override_model(self, model: str) -> None:
        self.current_config = SimpleNamespace(model=model)
        self.provider.model = model

    def resolve(self) -> _ArtifactProvider:
        return self.provider


class _ProviderSelector:
    def __init__(self, provider: _ArtifactProvider) -> None:
        self.provider = provider

    def clone(self) -> _SelectorClone:
        return _SelectorClone(self.provider)


class _FailedPublishProvider:
    provider_name = "test"

    def __init__(self) -> None:
        self.calls = 0
        self.model = "test/model"

    def chat(self, messages: list[Message], tools=None, config=None) -> AsyncIterator[Any]:
        self.calls += 1
        return self._stream(self.calls)

    async def _stream(self, call_number: int) -> AsyncIterator[Any]:
        if call_number == 1:
            yield ProviderToolUseStart(
                tool_use_id="publish-1",
                tool_name="publish_artifact",
            )
            yield ProviderToolUseEnd(
                tool_use_id="publish-1",
                tool_name="publish_artifact",
                arguments={"path": "missing-report.pptx"},
            )
            yield ProviderDone(stop_reason="tool_use", input_tokens=1, output_tokens=1)
            return
        yield ProviderText(text="Report file is ready for download.")
        yield ProviderDone(stop_reason="stop", input_tokens=1, output_tokens=1)

    async def list_models(self) -> list[ModelInfo]:
        return []


class _RetryPublishProvider(_FailedPublishProvider):
    async def _stream(self, call_number: int) -> AsyncIterator[Any]:
        yield ProviderText(text="Regenerating the presentation. ")
        yield ProviderToolUseStart(
            tool_use_id=f"publish-{call_number}",
            tool_name="publish_artifact",
        )
        yield ProviderToolUseEnd(
            tool_use_id=f"publish-{call_number}",
            tool_name="publish_artifact",
            arguments={"path": "report.pptx"},
        )
        yield ProviderDone(stop_reason="tool_use", input_tokens=1, output_tokens=1)


class _FailedCreatePptxProvider(_FailedPublishProvider):
    async def _stream(self, call_number: int) -> AsyncIterator[Any]:
        if call_number == 1:
            yield ProviderToolUseStart(tool_use_id="create-1", tool_name="create_pptx")
            yield ProviderToolUseEnd(
                tool_use_id="create-1",
                tool_name="create_pptx",
                arguments={"name": "report.pptx", "slides": [{"title": "Report"}]},
            )
            yield ProviderDone(stop_reason="tool_use", input_tokens=1, output_tokens=1)
            return
        yield ProviderText(text="Report file is ready for download.")
        yield ProviderDone(stop_reason="stop", input_tokens=1, output_tokens=1)


class _OmittedPublishProvider:
    provider_name = "test"

    def __init__(self) -> None:
        self.calls = 0
        self.model = "test/model"

    def chat(self, messages: list[Message], tools=None, config=None) -> AsyncIterator[Any]:
        self.calls += 1
        return self._stream(self.calls)

    async def _stream(self, call_number: int) -> AsyncIterator[Any]:
        if call_number == 1:
            yield ProviderToolUseStart(
                tool_use_id="write-1",
                tool_name="write_file",
            )
            yield ProviderToolUseEnd(
                tool_use_id="write-1",
                tool_name="write_file",
                arguments={
                    "path": "manual-big-write.html",
                    "content": "<!doctype html><title>Manual</title>",
                },
            )
            yield ProviderDone(stop_reason="tool_use", input_tokens=1, output_tokens=1)
            return
        yield ProviderText(text="Created manual-big-write.html for you.")
        yield ProviderDone(stop_reason="stop", input_tokens=1, output_tokens=1)

    async def list_models(self) -> list[ModelInfo]:
        return []


class _OmittedInvalidPptxProvider:
    provider_name = "test"

    def __init__(self) -> None:
        self.calls = 0
        self.model = "test/model"

    def chat(self, messages: list[Message], tools=None, config=None) -> AsyncIterator[Any]:
        self.calls += 1
        return self._stream(self.calls)

    async def _stream(self, call_number: int) -> AsyncIterator[Any]:
        if call_number == 1:
            yield ProviderToolUseStart(
                tool_use_id="write-1",
                tool_name="write_file",
            )
            yield ProviderToolUseEnd(
                tool_use_id="write-1",
                tool_name="write_file",
                arguments={
                    "path": "broken.pptx",
                    "content": "this is not a PowerPoint package",
                },
            )
            yield ProviderDone(stop_reason="tool_use", input_tokens=1, output_tokens=1)
            return
        yield ProviderText(text="Created broken.pptx for you.")
        yield ProviderDone(stop_reason="stop", input_tokens=1, output_tokens=1)

    async def list_models(self) -> list[ModelInfo]:
        return []


class _OmittedPatchPublishProvider:
    provider_name = "test"

    def __init__(self) -> None:
        self.calls = 0
        self.model = "test/model"

    def chat(self, messages: list[Message], tools=None, config=None) -> AsyncIterator[Any]:
        self.calls += 1
        return self._stream(self.calls)

    async def _stream(self, call_number: int) -> AsyncIterator[Any]:
        if call_number == 1:
            yield ProviderToolUseStart(
                tool_use_id="patch-1",
                tool_name="apply_patch",
            )
            yield ProviderToolUseEnd(
                tool_use_id="patch-1",
                tool_name="apply_patch",
                arguments={
                    "patch": (
                        "*** Begin Patch\n"
                        "*** Add File: patched.html\n"
                        "+<!doctype html><title>Patched</title>\n"
                        "*** End Patch\n"
                    ),
                },
            )
            yield ProviderDone(stop_reason="tool_use", input_tokens=1, output_tokens=1)
            return
        yield ProviderText(text="Created patched.html.")
        yield ProviderDone(stop_reason="stop", input_tokens=1, output_tokens=1)

    async def list_models(self) -> list[ModelInfo]:
        return []


class _EditedConfigProvider:
    provider_name = "test"

    def __init__(self) -> None:
        self.calls = 0
        self.model = "test/model"

    def chat(self, messages: list[Message], tools=None, config=None) -> AsyncIterator[Any]:
        self.calls += 1
        return self._stream(self.calls)

    async def _stream(self, call_number: int) -> AsyncIterator[Any]:
        if call_number == 1:
            yield ProviderToolUseStart(
                tool_use_id="edit-1",
                tool_name="edit_file",
            )
            yield ProviderToolUseEnd(
                tool_use_id="edit-1",
                tool_name="edit_file",
                arguments={
                    "path": "config.json",
                    "old_text": "\"enabled\": false",
                    "new_text": "\"enabled\": true",
                },
            )
            yield ProviderDone(stop_reason="tool_use", input_tokens=1, output_tokens=1)
            return
        yield ProviderText(text="Updated config.json.")
        yield ProviderDone(stop_reason="stop", input_tokens=1, output_tokens=1)

    async def list_models(self) -> list[ModelInfo]:
        return []


class _MixedSizeOmittedPublishProvider:
    provider_name = "test"

    def __init__(self) -> None:
        self.calls = 0
        self.model = "test/model"

    def chat(self, messages: list[Message], tools=None, config=None) -> AsyncIterator[Any]:
        self.calls += 1
        return self._stream(self.calls)

    async def _stream(self, call_number: int) -> AsyncIterator[Any]:
        if call_number == 1:
            for index, payload in enumerate(
                (
                    {"path": "small.html", "content": "<title>ok</title>"},
                    {"path": "large.html", "content": "<title>" + ("x" * 80) + "</title>"},
                ),
                start=1,
            ):
                yield ProviderToolUseStart(
                    tool_use_id=f"write-{index}",
                    tool_name="write_file",
                )
                yield ProviderToolUseEnd(
                    tool_use_id=f"write-{index}",
                    tool_name="write_file",
                    arguments=payload,
                )
            yield ProviderDone(stop_reason="tool_use", input_tokens=1, output_tokens=1)
            return
        yield ProviderText(text="Created small.html and large.html.")
        yield ProviderDone(stop_reason="stop", input_tokens=1, output_tokens=1)

    async def list_models(self) -> list[ModelInfo]:
        return []


class _MemoryJsonWriteProvider:
    provider_name = "test"

    def __init__(self) -> None:
        self.calls = 0
        self.model = "test/model"

    def chat(self, messages: list[Message], tools=None, config=None) -> AsyncIterator[Any]:
        self.calls += 1
        return self._stream(self.calls)

    async def _stream(self, call_number: int) -> AsyncIterator[Any]:
        if call_number == 1:
            yield ProviderToolUseStart(
                tool_use_id="write-1",
                tool_name="write_file",
            )
            yield ProviderToolUseEnd(
                tool_use_id="write-1",
                tool_name="write_file",
                arguments={
                    "path": "memory/cache.json",
                    "content": "{\"state\":\"internal\"}",
                },
            )
            yield ProviderDone(stop_reason="tool_use", input_tokens=1, output_tokens=1)
            return
        yield ProviderText(text="Updated memory/cache.json.")
        yield ProviderDone(stop_reason="stop", input_tokens=1, output_tokens=1)

    async def list_models(self) -> list[ModelInfo]:
        return []


class _SameContentOmittedPublishProvider:
    provider_name = "test"

    def __init__(self) -> None:
        self.calls = 0
        self.model = "test/model"

    def chat(self, messages: list[Message], tools=None, config=None) -> AsyncIterator[Any]:
        self.calls += 1
        return self._stream(self.calls)

    async def _stream(self, call_number: int) -> AsyncIterator[Any]:
        if call_number == 1:
            for index, path in enumerate(("first.html", "second.html"), start=1):
                yield ProviderToolUseStart(
                    tool_use_id=f"write-{index}",
                    tool_name="write_file",
                )
                yield ProviderToolUseEnd(
                    tool_use_id=f"write-{index}",
                    tool_name="write_file",
                    arguments={
                        "path": path,
                        "content": "<!doctype html><title>Same</title>",
                    },
                )
            yield ProviderDone(stop_reason="tool_use", input_tokens=1, output_tokens=1)
            return
        yield ProviderText(text="Created first.html and second.html.")
        yield ProviderDone(stop_reason="stop", input_tokens=1, output_tokens=1)

    async def list_models(self) -> list[ModelInfo]:
        return []


class _PartialOmittedPublishProvider:
    provider_name = "test"

    def __init__(self) -> None:
        self.calls = 0
        self.model = "test/model"

    def chat(self, messages: list[Message], tools=None, config=None) -> AsyncIterator[Any]:
        self.calls += 1
        return self._stream(self.calls)

    async def _stream(self, call_number: int) -> AsyncIterator[Any]:
        if call_number == 1:
            yield ProviderToolUseStart(tool_use_id="make-1", tool_name="make_file")
            yield ProviderToolUseEnd(
                tool_use_id="make-1",
                tool_name="make_file",
                arguments={},
            )
            yield ProviderToolUseStart(
                tool_use_id="write-1",
                tool_name="write_file",
            )
            yield ProviderToolUseEnd(
                tool_use_id="write-1",
                tool_name="write_file",
                arguments={
                    "path": "second.html",
                    "content": "<!doctype html><title>Second</title>",
                },
            )
            yield ProviderDone(stop_reason="tool_use", input_tokens=1, output_tokens=1)
            return
        yield ProviderText(text="Created runtime.txt and second.html.")
        yield ProviderDone(stop_reason="stop", input_tokens=1, output_tokens=1)

    async def list_models(self) -> list[ModelInfo]:
        return []


def _registry() -> ToolRegistry:
    registry = ToolRegistry()

    async def make_file() -> str:
        ctx = current_tool_context.get()
        assert ctx is not None
        ctx.published_artifacts.append(
            {
                "id": "art-runtime",
                "kind": "artifact_ref",
                "name": "runtime.txt",
                "mime": "text/plain",
                "size": 4,
                "sha256": "b" * 64,
                "session_id": ctx.artifact_session_id,
                "session_key": ctx.session_key,
                "source": "make_file",
                "created_at": "2026-05-06T12:00:00Z",
                "download_url": (
                    "/api/v1/artifacts/art-runtime"
                    "?sessionKey=agent%3Amain%3Awebchat%3Aartifact-runtime"
                ),
            }
        )
        return "published"

    registry.register(
        ToolSpec(name="make_file", description="Make a file", parameters={}),
        make_file,
    )
    return registry


def _registry_with_write_file() -> ToolRegistry:
    registry = _registry()
    write_file = filesystem.write_file.__wrapped__.__wrapped__  # type: ignore[attr-defined]
    registry.register(
        ToolSpec(
            name="write_file",
            description="Write a file",
            parameters={
                "type": "object",
                "properties": {
                    "path": {"type": "string"},
                    "content": {"type": "string"},
                },
                "required": ["path", "content"],
            },
        ),
        write_file,
    )
    return registry


def _apply_patch_registry() -> ToolRegistry:
    registry = ToolRegistry()
    apply_patch = patch_tools.apply_patch.__wrapped__.__wrapped__  # type: ignore[attr-defined]
    registry.register(
        ToolSpec(
            name="apply_patch",
            description="Apply a patch",
            parameters={
                "type": "object",
                "properties": {
                    "patch": {"type": "string"},
                },
                "required": ["patch"],
            },
        ),
        apply_patch,
    )
    return registry


def _edit_file_registry() -> ToolRegistry:
    registry = ToolRegistry()
    edit_file = filesystem.edit_file.__wrapped__.__wrapped__  # type: ignore[attr-defined]
    registry.register(
        ToolSpec(
            name="edit_file",
            description="Edit a file",
            parameters={
                "type": "object",
                "properties": {
                    "path": {"type": "string"},
                    "old_text": {"type": "string"},
                    "new_text": {"type": "string"},
                },
                "required": ["path", "old_text", "new_text"],
            },
        ),
        edit_file,
    )
    return registry


def _write_file_registry() -> ToolRegistry:
    registry = ToolRegistry()
    write_file = filesystem.write_file.__wrapped__.__wrapped__  # type: ignore[attr-defined]
    registry.register(
        ToolSpec(
            name="write_file",
            description="Write a file",
            parameters={
                "type": "object",
                "properties": {
                    "path": {"type": "string"},
                    "content": {"type": "string"},
                },
                "required": ["path", "content"],
            },
        ),
        write_file,
    )
    return registry


def _failed_publish_registry() -> ToolRegistry:
    registry = ToolRegistry()

    async def publish_artifact(path: str) -> str:
        raise ToolError(f"artifact file not found: {path}")

    registry.register(
        ToolSpec(
            name="publish_artifact",
            description="Publish a generated artifact",
            parameters={
                "type": "object",
                "properties": {"path": {"type": "string"}},
                "required": ["path"],
            },
        ),
        publish_artifact,
    )
    return registry


def _retry_publish_registry() -> ToolRegistry:
    registry = ToolRegistry()
    calls = 0

    async def publish_artifact(path: str) -> str:
        nonlocal calls
        calls += 1
        if calls == 1:
            raise RetryableToolInputError("Regenerate the invalid PPTX and try again.")
        ctx = current_tool_context.get()
        assert ctx is not None
        ctx.published_artifacts.append(
            {
                "id": "art-retried",
                "kind": "artifact_ref",
                "name": path,
                "mime": (
                    "application/vnd.openxmlformats-officedocument."
                    "presentationml.presentation"
                ),
                "size": 8,
                "sha256": "d" * 64,
                "session_id": ctx.artifact_session_id,
                "session_key": ctx.session_key,
                "source": "publish_artifact",
                "created_at": "2026-07-20T00:00:00Z",
                "download_url": "/api/v1/artifacts/art-retried",
            }
        )
        return json.dumps({"status": "published", "artifact": {"name": path}})

    registry.register(
        ToolSpec(
            name="publish_artifact",
            description="Publish a generated artifact",
            parameters={
                "type": "object",
                "properties": {"path": {"type": "string"}},
                "required": ["path"],
            },
        ),
        publish_artifact,
    )
    return registry


def _failed_create_pptx_registry() -> ToolRegistry:
    registry = ToolRegistry()

    async def create_pptx(slides: list[dict[str, Any]], name: str | None = None) -> str:
        raise RetryableToolInputError("The PPTX was not attached; regenerate it.")

    registry.register(
        ToolSpec(
            name="create_pptx",
            description="Create a presentation",
            parameters={
                "type": "object",
                "properties": {
                    "name": {"type": "string"},
                    "slides": {"type": "array", "items": {"type": "object"}},
                },
                "required": ["slides"],
            },
        ),
        create_pptx,
    )
    return registry


def _publish_then_forbidden_tool_registry() -> tuple[ToolRegistry, list[str]]:
    registry = ToolRegistry()
    forbidden_calls: list[str] = []

    async def publish_artifact(path: str) -> str:
        ctx = current_tool_context.get()
        assert ctx is not None
        ctx.published_artifacts.append(
            {
                "id": "art-published",
                "kind": "artifact_ref",
                "name": path,
                "mime": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                "size": 8,
                "sha256": "c" * 64,
                "session_id": ctx.artifact_session_id,
                "session_key": ctx.session_key,
                "source": "publish_artifact",
                "created_at": "2026-05-06T12:00:00Z",
                "download_url": "/api/v1/artifacts/art-published",
            }
        )
        return json.dumps({"status": "published", "artifact": {"name": path}})

    async def qa_check(path: str) -> str:
        forbidden_calls.append(path)
        return "qa done"

    registry.register(
        ToolSpec(
            name="publish_artifact",
            description="Publish a generated artifact",
            parameters={
                "type": "object",
                "properties": {"path": {"type": "string"}},
                "required": ["path"],
            },
        ),
        publish_artifact,
    )
    registry.register(
        ToolSpec(
            name="qa_check",
            description="QA check",
            parameters={
                "type": "object",
                "properties": {"path": {"type": "string"}},
                "required": ["path"],
            },
        ),
        qa_check,
    )
    return registry, forbidden_calls


def _goal_publish_loop_registry(
    *,
    qa_fails: bool = False,
) -> tuple[
    ToolRegistry,
    list[str],
    list[str],
]:
    registry = ToolRegistry()
    control_calls: list[str] = []
    qa_calls: list[str] = []

    async def publish_artifact(path: str) -> str:
        ctx = current_tool_context.get()
        assert ctx is not None
        ctx.published_artifacts.append(
            {
                "id": "art-goal-published",
                "kind": "artifact_ref",
                "name": path,
                "mime": "text/html",
                "size": 8,
                "sha256": "e" * 64,
                "session_id": ctx.artifact_session_id,
                "session_key": ctx.session_key,
                "source": "publish_artifact",
                "created_at": "2026-08-08T00:00:00Z",
                "download_url": "/api/v1/artifacts/art-goal-published",
            }
        )
        return json.dumps({"status": "published", "artifact": {"name": path}})

    async def update_goal_progress(steps: list[dict[str, Any]]) -> str:
        control_calls.append(f"progress:{steps[0]['status']}")
        return json.dumps({"status": "accepted"})

    async def update_goal(status: str, reason: str | None = None) -> str:
        control_calls.append(f"goal:{status}")
        return json.dumps({"status": "accepted", "goal": {"status": status}})

    async def qa_check(path: str) -> str:
        qa_calls.append(path)
        if qa_fails:
            raise ToolError("synthetic QA failure")
        return "qa done"

    registry.register(
        ToolSpec(
            name="publish_artifact",
            description="Publish a generated artifact",
            parameters={
                "type": "object",
                "properties": {"path": {"type": "string"}},
                "required": ["path"],
            },
        ),
        publish_artifact,
    )
    registry.register(
        ToolSpec(
            name="update_goal_progress",
            description="Update Goal progress",
            parameters={
                "type": "object",
                "properties": {"steps": {"type": "array"}},
                "required": ["steps"],
            },
            exposed_by_default=False,
        ),
        update_goal_progress,
    )
    registry.register(
        ToolSpec(
            name="update_goal",
            description="Update Goal status",
            parameters={
                "type": "object",
                "properties": {
                    "status": {"type": "string"},
                    "reason": {"type": "string"},
                },
                "required": ["status"],
            },
            exposed_by_default=False,
        ),
        update_goal,
    )
    registry.register(
        ToolSpec(
            name="qa_check",
            description="QA check",
            parameters={
                "type": "object",
                "properties": {"path": {"type": "string"}},
                "required": ["path"],
            },
        ),
        qa_check,
    )
    return registry, control_calls, qa_calls


@pytest.mark.asyncio
async def test_turn_runner_streams_artifact_event_and_persists_history(tmp_path) -> None:
    storage = SessionStorage(":memory:")
    await storage.connect()
    manager = SessionManager(storage)
    session_key = "agent:main:webchat:artifact-runtime"
    session = await manager.create(session_key)
    runner = TurnRunner(
        provider_selector=_ProviderSelector(_ArtifactProvider()),
        tool_registry=_registry(),
        session_manager=manager,
        config=GatewayConfig(
            attachments=AttachmentsConfig(media_root=str(tmp_path / "media")),
            squilla_router=SquillaRouterConfig(enabled=False),
        ),
    )
    tool_context = ToolContext(
        is_owner=True,
        caller_kind=CallerKind.WEB,
        workspace_dir=str(tmp_path),
    )

    try:
        events = [
            event
            async for event in runner.run(
                "make it",
                session_key,
                tool_context=tool_context,
                history_has_persisted_user=False,
                no_memory_capture=True,
            )
        ]
        artifact_events = [event for event in events if isinstance(event, ArtifactEvent)]
        assert len(artifact_events) == 1
        assert artifact_events[0].id == "art-runtime"
        assert artifact_events[0].session_id == session.session_id
        assert artifact_events[0].session_key == ""
        assert artifact_events[0].download_url == "/api/v1/artifacts/art-runtime"

        transcript = await manager.get_transcript(session_key)
        assistant = [entry for entry in transcript if entry.role == "assistant"][-1]
        payload = json.loads(assistant.content)
        assert payload["text"] == "done"
        assert payload["artifacts"][0]["id"] == "art-runtime"
        assert payload["artifacts"][0]["session_id"] == session.session_id
        assert "session_key" not in payload["artifacts"][0]
        assert "sessionKey" not in assistant.content

        class _HistoryCapture:
            def __init__(self) -> None:
                self.history = []

            def set_history(self, history) -> None:
                self.history = history

        history_capture = _HistoryCapture()
        await runner._load_history(agent=history_capture, session_key=session_key)
        assert "[generated artifact omitted: runtime.txt (text/plain)]" in str(
            history_capture.history[-1].content
        )
    finally:
        await storage.close()


@pytest.mark.asyncio
async def test_turn_runner_cancel_after_artifact_persists_recoverable_delivery_text(
    tmp_path,
) -> None:
    storage = SessionStorage(":memory:")
    await storage.connect()
    manager = SessionManager(storage)
    session_key = "agent:main:webchat:artifact-cancelled"
    session = await manager.create(session_key)
    runner = TurnRunner(
        provider_selector=_ProviderSelector(_BlockingAfterArtifactProvider()),
        tool_registry=_registry(),
        session_manager=manager,
        config=GatewayConfig(
            attachments=AttachmentsConfig(media_root=str(tmp_path / "media")),
            squilla_router=SquillaRouterConfig(enabled=False),
        ),
    )
    tool_context = ToolContext(
        is_owner=True,
        caller_kind=CallerKind.WEB,
        workspace_dir=str(tmp_path),
    )
    artifact_seen = asyncio.Event()

    async def _consume() -> None:
        async for event in runner.run(
            "make it",
            session_key,
            tool_context=tool_context,
            history_has_persisted_user=False,
            no_memory_capture=True,
        ):
            if isinstance(event, ArtifactEvent):
                artifact_seen.set()

    task = asyncio.create_task(_consume())
    try:
        await asyncio.wait_for(artifact_seen.wait(), timeout=2.0)
        await asyncio.sleep(0)
        task.cancel()
        with pytest.raises(asyncio.CancelledError):
            await task

        transcript = await manager.get_transcript(session_key)
        assistant = [entry for entry in transcript if entry.role == "assistant"][-1]
        payload = json.loads(assistant.content)
        assert payload["artifacts"][0]["id"] == "art-runtime"
        assert payload["artifacts"][0]["session_id"] == session.session_id
        assert "The generated file was delivered" in payload["text"]
        assert "[interrupted]" not in payload["text"]
    finally:
        if not task.done():
            task.cancel()
            with contextlib.suppress(asyncio.CancelledError):
                await task
        await storage.close()


@pytest.mark.asyncio
async def test_turn_runner_suppresses_tools_after_successful_publish_artifact(
    tmp_path,
) -> None:
    storage = SessionStorage(":memory:")
    await storage.connect()
    manager = SessionManager(storage)
    session_key = "agent:main:webchat:artifact-stop"
    session = await manager.create(session_key)
    provider = _PostPublishToolLoopProvider()
    registry, forbidden_calls = _publish_then_forbidden_tool_registry()
    runner = TurnRunner(
        provider_selector=_ProviderSelector(provider),
        tool_registry=registry,
        session_manager=manager,
        config=GatewayConfig(
            attachments=AttachmentsConfig(media_root=str(tmp_path / "media")),
            squilla_router=SquillaRouterConfig(enabled=False),
        ),
    )
    tool_context = ToolContext(
        is_owner=True,
        caller_kind=CallerKind.WEB,
        workspace_dir=str(tmp_path),
    )

    try:
        events = [
            event
            async for event in runner.run(
                "make ppt",
                session_key,
                tool_context=tool_context,
                history_has_persisted_user=False,
                no_memory_capture=True,
            )
        ]

        done = next(event for event in events if isinstance(event, DoneEvent))
        artifact_events = [event for event in events if isinstance(event, ArtifactEvent)]
        tool_starts = [event for event in events if isinstance(event, ToolUseStartEvent)]

        assert provider.calls == 1
        assert provider.tools_seen == [True]
        assert forbidden_calls == []
        assert [event.tool_name for event in tool_starts] == ["publish_artifact"]
        assert artifact_events[0].id == "art-published"
        assert artifact_events[0].session_id == session.session_id
        text_deltas = [event.text for event in events if isinstance(event, TextDeltaEvent)]
        assert "".join(text_deltas) == done.text
        assert done.text.startswith("Preparing your presentation.")
        assert "The generated file is ready" in done.text

        transcript = await manager.get_transcript(session_key)
        assistant = [entry for entry in transcript if entry.role == "assistant"][-1]
        payload = json.loads(assistant.content)
        assert payload["artifacts"][0]["id"] == "art-published"
        assert "The generated file is ready" in payload["text"]
    finally:
        await storage.close()


async def _run_goal_publish_loop(
    tmp_path,
    *,
    plain_final: bool,
    post_publish_error_code: str | None = None,
    post_publish_tool_failure: bool = False,
    terminal_status: str = "complete",
    illegal_terminal_summary_tool: bool = False,
    terminal_summary_mode: str = "normal",
    max_iterations: int | None = None,
) -> tuple[
    _GoalPostPublishLoopProvider,
    list[str],
    list[str],
    list[Any],
]:
    storage = SessionStorage(":memory:")
    await storage.connect()
    manager = SessionManager(storage)
    session_key = (
        "agent:main:webchat:goal-artifact-plain"
        if plain_final
        else "agent:main:webchat:goal-artifact-complete"
    )
    await manager.create(session_key)
    provider = _GoalPostPublishLoopProvider(
        plain_final=plain_final,
        post_publish_error_code=post_publish_error_code,
        post_publish_tool_failure=post_publish_tool_failure,
        terminal_status=terminal_status,
        illegal_terminal_summary_tool=illegal_terminal_summary_tool,
        terminal_summary_mode=terminal_summary_mode,
    )
    registry, control_calls, qa_calls = _goal_publish_loop_registry(
        qa_fails=post_publish_tool_failure,
    )
    runner = TurnRunner(
        provider_selector=_ProviderSelector(provider),
        tool_registry=registry,
        session_manager=manager,
        config=GatewayConfig(
            attachments=AttachmentsConfig(media_root=str(tmp_path / "media")),
            squilla_router=SquillaRouterConfig(enabled=False),
            agent_max_provider_retries=0,
        ),
    )
    tool_context = ToolContext(
        is_owner=True,
        caller_kind=CallerKind.WEB,
        workspace_dir=str(tmp_path),
        session_key=session_key,
        task_id="task-goal-artifact",
        collaboration_mode="default",
        goal_context={
            "schemaVersion": 1,
            "sessionId": "session-goal-artifact",
            "epoch": 0,
            "goalId": "goal-artifact",
            "objectiveRevision": 1,
            "objectiveSnapshot": "Publish and verify the report.",
            "taskId": "task-goal-artifact",
            "continuationSeq": 0,
            "automatic": False,
            "progress": {
                "steps": [
                    {"step": "Publish the verified artifact", "status": "in_progress"}
                ]
            },
        },
    )
    try:
        events = [
            event
            async for event in runner.run(
                "publish the report",
                session_key,
                tool_context=tool_context,
                history_has_persisted_user=False,
                no_memory_capture=True,
                max_iterations=max_iterations,
            )
        ]
        return provider, control_calls, qa_calls, events
    finally:
        await storage.close()


def _assert_goal_artifact_published_once(events: list[Any]) -> None:
    assert len([event for event in events if isinstance(event, ArtifactEvent)]) == 1
    assert [
        event.tool_name
        for event in events
        if isinstance(event, ToolUseStartEvent)
        and event.tool_name == "publish_artifact"
    ] == ["publish_artifact"]


@pytest.mark.asyncio
async def test_goal_publish_continues_normal_loop_through_terminal_and_final_summary(
    tmp_path,
) -> None:
    provider, control_calls, qa_calls, events = await _run_goal_publish_loop(
        tmp_path,
        plain_final=False,
    )

    assert provider.calls == 4
    expected_tools = {
        "publish_artifact",
        "qa_check",
        "update_goal",
        "update_goal_progress",
    }
    assert all(set(tool_names) == expected_tools for tool_names in provider.tool_names_seen[:3])
    assert provider.tool_names_seen[3] == []
    assert isinstance(provider.requests[3][-1].content, list)
    assert any(
        isinstance(block, ContentBlockToolResult)
        for block in provider.requests[3][-1].content
    )
    assert not any(
        isinstance(message.content, str)
        and message.content.startswith(
            ("[Runtime ", "Progress check:", "Time check:", "STOP: multiple")
        )
        for message in provider.requests[3]
    )
    assert control_calls == ["progress:completed", "goal:complete"]
    assert qa_calls == []
    _assert_goal_artifact_published_once(events)
    assert not any(isinstance(event, ErrorEvent) for event in events)
    assert "Goal delivery checkpoint" not in "\n".join(
        str(message.content)
        for request in provider.requests
        for message in request
    )
    assert [
        event.tool_name
        for event in events
        if isinstance(event, ToolUseStartEvent)
    ] == ["publish_artifact", "update_goal_progress", "update_goal"]
    done = next(event for event in events if isinstance(event, DoneEvent))
    assert done.text == "The Goal is complete."


@pytest.mark.asyncio
@pytest.mark.parametrize("terminal_status", ["complete", "blocked"])
async def test_goal_terminal_summary_disables_and_ignores_provider_tool_calls(
    tmp_path,
    terminal_status: str,
) -> None:
    provider, control_calls, qa_calls, events = await _run_goal_publish_loop(
        tmp_path,
        plain_final=False,
        terminal_status=terminal_status,
        illegal_terminal_summary_tool=True,
    )

    assert provider.calls == 4
    assert provider.tool_names_seen[3] == []
    assert control_calls == ["progress:completed", f"goal:{terminal_status}"]
    assert qa_calls == []
    _assert_goal_artifact_published_once(events)
    assert not any(isinstance(event, ErrorEvent) for event in events)
    assert not any(
        isinstance(event, ToolUseStartEvent) and event.tool_name == "qa_check"
        for event in events
    )
    done = next(event for event in events if isinstance(event, DoneEvent))
    assert done.text == f"The Goal is {terminal_status}."


@pytest.mark.asyncio
@pytest.mark.parametrize("terminal_status", ["complete", "blocked"])
async def test_goal_terminal_on_last_iteration_uses_terminal_summary_not_partial(
    tmp_path,
    terminal_status: str,
) -> None:
    provider, control_calls, qa_calls, events = await _run_goal_publish_loop(
        tmp_path,
        plain_final=False,
        terminal_status=terminal_status,
        illegal_terminal_summary_tool=True,
        max_iterations=3,
    )

    assert provider.calls == 4
    assert provider.tool_names_seen[3] == []
    final_request_text = "\n".join(
        str(message.content) for message in provider.requests[3]
    )
    assert "The configured iteration limit has been reached" not in final_request_text
    assert "best concise final answer from the work completed so far" not in final_request_text
    assert control_calls == ["progress:completed", f"goal:{terminal_status}"]
    assert qa_calls == []
    _assert_goal_artifact_published_once(events)
    assert not any(isinstance(event, ErrorEvent) for event in events)
    done = next(event for event in events if isinstance(event, DoneEvent))
    assert done.text == f"The Goal is {terminal_status}."


@pytest.mark.asyncio
@pytest.mark.parametrize(
    "terminal_summary_mode",
    ["empty", "reasoning_only", "stream_incomplete", "length_capped"],
)
async def test_goal_terminal_invalid_summary_degrades_without_system_error(
    tmp_path,
    terminal_summary_mode: str,
) -> None:
    provider, control_calls, qa_calls, events = await _run_goal_publish_loop(
        tmp_path,
        plain_final=False,
        terminal_summary_mode=terminal_summary_mode,
    )

    assert provider.calls == 4
    assert control_calls == ["progress:completed", "goal:complete"]
    assert qa_calls == []
    assert not any(isinstance(event, ErrorEvent) for event in events)
    done = next(event for event in events if isinstance(event, DoneEvent))
    assert done.text.endswith("The Goal is complete.")


@pytest.mark.asyncio
@pytest.mark.parametrize(
    "preflight_failure",
    ["request_assembly", "request_validation", "identical_request_abort"],
)
@pytest.mark.parametrize("terminal_status", ["complete", "blocked"])
async def test_goal_terminal_summary_preflight_failure_degrades_without_system_error(
    tmp_path,
    monkeypatch: pytest.MonkeyPatch,
    preflight_failure: str,
    terminal_status: str,
) -> None:
    if preflight_failure == "request_assembly":
        original_assemble = Agent._provider_request_messages_with_sanitize_async

        async def fail_terminal_summary_assembly(
            agent: Agent,
            *args: Any,
            **kwargs: Any,
        ) -> Any:
            if getattr(agent.provider, "calls", 0) == 3:
                raise OSError("synthetic terminal summary assembly failure")
            return await original_assemble(agent, *args, **kwargs)

        monkeypatch.setattr(
            Agent,
            "_provider_request_messages_with_sanitize_async",
            fail_terminal_summary_assembly,
        )
    elif preflight_failure == "request_validation":
        original_validate = agent_module.validate_provider_chat_request

        def fail_terminal_summary_validation(provider: Any, messages: list[Message]) -> Any:
            if getattr(provider, "calls", 0) == 3:
                return ProviderError(
                    message="Synthetic terminal summary validation failure.",
                    code="synthetic_terminal_summary_validation_failure",
                )
            return original_validate(provider, messages)

        monkeypatch.setattr(
            agent_module,
            "validate_provider_chat_request",
            fail_terminal_summary_validation,
        )
    else:
        original_identical_action = Agent._identical_request_loop_break_action

        def abort_terminal_summary_request(
            agent: Agent,
            request_messages: list[Message],
            *,
            first_attempt: bool,
        ) -> str | None:
            if getattr(agent.provider, "calls", 0) == 3:
                return "abort"
            return original_identical_action(
                agent,
                request_messages,
                first_attempt=first_attempt,
            )

        monkeypatch.setattr(
            Agent,
            "_identical_request_loop_break_action",
            abort_terminal_summary_request,
        )

    provider, control_calls, qa_calls, events = await _run_goal_publish_loop(
        tmp_path,
        plain_final=False,
        terminal_status=terminal_status,
    )

    assert provider.calls == 3
    assert control_calls == ["progress:completed", f"goal:{terminal_status}"]
    assert qa_calls == []
    assert not any(isinstance(event, ErrorEvent) for event in events)
    done = next(event for event in events if isinstance(event, DoneEvent))
    assert done.text == f"The Goal is {terminal_status}."


@pytest.mark.asyncio
async def test_goal_publish_then_plain_final_succeeds_without_artificial_error(
    tmp_path,
) -> None:
    provider, control_calls, qa_calls, events = await _run_goal_publish_loop(
        tmp_path,
        plain_final=True,
    )

    assert provider.calls == 2
    assert set(provider.tool_names_seen[1]) == {
        "publish_artifact",
        "qa_check",
        "update_goal",
        "update_goal_progress",
    }
    assert control_calls == []
    assert qa_calls == []
    _assert_goal_artifact_published_once(events)
    assert not any(isinstance(event, ErrorEvent) for event in events)
    done = next(event for event in events if isinstance(event, DoneEvent))
    assert done.text == "The report is ready."


@pytest.mark.asyncio
@pytest.mark.parametrize(
    "provider_error_code",
    ["synthetic_upstream_failure", "usage_limit_reached"],
)
async def test_goal_publish_then_provider_failure_uses_normal_error_contract(
    tmp_path,
    provider_error_code: str,
) -> None:
    provider, control_calls, qa_calls, events = await _run_goal_publish_loop(
        tmp_path,
        plain_final=False,
        post_publish_error_code=provider_error_code,
    )

    assert provider.calls == 2
    assert control_calls == []
    assert qa_calls == []
    _assert_goal_artifact_published_once(events)
    assert [
        event.code for event in events if isinstance(event, ErrorEvent)
    ] == [provider_error_code]
    # TurnRunner always emits one terminal Done envelope for accounting, even
    # for failed turns. It must not convert this real provider failure into the
    # deterministic non-Goal artifact-success summary.
    done = next(event for event in events if isinstance(event, DoneEvent))
    assert done.text == ""


@pytest.mark.asyncio
async def test_goal_publish_then_tool_failure_returns_to_normal_provider_loop(
    tmp_path,
) -> None:
    provider, control_calls, qa_calls, events = await _run_goal_publish_loop(
        tmp_path,
        plain_final=False,
        post_publish_tool_failure=True,
    )

    assert provider.calls == 3
    assert control_calls == []
    assert qa_calls == ["report.html"]
    _assert_goal_artifact_published_once(events)
    qa_results = [
        event
        for event in events
        if isinstance(event, ToolResultEvent) and event.tool_name == "qa_check"
    ]
    assert len(qa_results) == 1
    assert qa_results[0].is_error is True
    qa_error = json.loads(str(qa_results[0].result))
    assert qa_error["status"] == "error"
    assert qa_error["error_class"] == "ToolError"
    assert not any(isinstance(event, ErrorEvent) for event in events)
    done = next(event for event in events if isinstance(event, DoneEvent))
    assert done.text == "The published report remains available for review."


@pytest.mark.asyncio
async def test_goal_terminal_batch_skips_sessions_yield_then_summarizes(
    tmp_path,
) -> None:
    # Exercise the shared Agent state machine directly: surface policy decides
    # who may call sessions_yield, while this test locks the result precedence.
    provider = _GoalPublishAndYieldProvider()

    executed: list[str] = []

    async def handle_tool(call: ToolCall) -> ToolResult:
        executed.append(call.tool_name)
        if call.tool_name == "publish_artifact":
            return ToolResult(
                tool_use_id=call.tool_use_id,
                tool_name=call.tool_name,
                content=json.dumps({"status": "published"}),
                artifacts=[
                    {
                        "id": "art-goal-yield",
                        "kind": "artifact_ref",
                        "name": "report.html",
                        "mime": "text/html",
                        "size": 8,
                        "sha256": "f" * 64,
                        "session_id": "session-goal-artifact-yield",
                        "session_key": "agent:main:webchat:goal-artifact-yield",
                        "source": "publish_artifact",
                        "created_at": "2026-08-08T00:00:00Z",
                        "download_url": "/api/v1/artifacts/art-goal-yield",
                    }
                ],
            )
        if call.tool_name == "update_goal":
            return ToolResult(
                tool_use_id=call.tool_use_id,
                tool_name=call.tool_name,
                content=json.dumps(
                    {"status": "accepted", "goal": {"status": "complete"}}
                ),
            )
        assert call.tool_name == "sessions_yield"
        return ToolResult(
            tool_use_id=call.tool_use_id,
            tool_name=call.tool_name,
            content=json.dumps({"status": "yielded"}),
        )

    tool_context = ToolContext(
        is_owner=True,
        caller_kind=CallerKind.WEB,
        workspace_dir=str(tmp_path),
        session_key="agent:main:webchat:goal-artifact-yield",
        task_id="task-goal-artifact-yield",
        collaboration_mode="default",
        goal_context={
            "schemaVersion": 1,
            "sessionId": "session-goal-artifact-yield",
            "epoch": 0,
            "goalId": "goal-artifact-yield",
            "objectiveRevision": 1,
            "objectiveSnapshot": "Publish the report, then yield to the scheduler.",
            "taskId": "task-goal-artifact-yield",
            "continuationSeq": 0,
            "automatic": False,
            "progress": {"steps": []},
        },
    )
    agent = Agent(
        provider=provider,
        config=AgentConfig(max_provider_retries=0),
        tool_definitions=[
            ToolDefinition(
                name="publish_artifact",
                description="Publish a generated artifact",
                input_schema=ToolInputSchema(),
            ),
            ToolDefinition(
                name="update_goal",
                description="Complete the Goal",
                input_schema=ToolInputSchema(),
            ),
            ToolDefinition(
                name="sessions_yield",
                description="Yield the current turn",
                input_schema=ToolInputSchema(),
            ),
        ],
        tool_handler=handle_tool,
        tool_context=tool_context,
    )

    events = [event async for event in agent.run_turn("publish and yield")]

    assert provider.calls == 2
    assert provider.tools_by_call[1] == []
    assert executed == ["publish_artifact", "update_goal"]
    _assert_goal_artifact_published_once(events)
    assert [
        event.tool_name
        for event in events
        if isinstance(event, ToolUseStartEvent)
    ] == ["publish_artifact", "update_goal", "sessions_yield"]
    yield_result = next(
        event
        for event in events
        if isinstance(event, ToolResultEvent) and event.tool_name == "sessions_yield"
    )
    assert json.loads(str(yield_result.result)) == {
        "status": "not_executed",
        "reason": "prior_tool_dispatch_boundary",
        "boundary_tool": "update_goal",
        "boundary_tool_use_id": "goal-1",
    }
    assert yield_result.is_error is True
    assert not any(isinstance(event, ErrorEvent) for event in events)
    assert "The Goal is complete and the artifact is ready." in "".join(
        event.text for event in events if isinstance(event, TextDeltaEvent)
    )


@pytest.mark.asyncio
@pytest.mark.parametrize(
    ("terminal_accepted", "max_turn_llm_calls", "summary_mode"),
    [
        (True, 0, "normal"),
        (True, 1, "normal"),
        (False, 0, "normal"),
        (True, 0, "reasoning_stream"),
        (True, 0, "thinking_error"),
    ],
)
async def test_goal_terminal_result_is_an_immediate_tool_dispatch_boundary(
    tmp_path,
    terminal_accepted: bool,
    max_turn_llm_calls: int,
    summary_mode: str,
) -> None:
    provider = _GoalTerminalThenPublishProvider(summary_mode=summary_mode)
    executed: list[str] = []

    async def handle_tool(call: ToolCall) -> ToolResult:
        executed.append(call.tool_name)
        if call.tool_name == "update_goal":
            return ToolResult(
                tool_use_id=call.tool_use_id,
                tool_name=call.tool_name,
                content=json.dumps(
                    {"status": "accepted", "goal": {"status": "complete"}}
                    if terminal_accepted
                    else {"status": "rejected", "code": "STALE_GOAL"}
                ),
                is_error=not terminal_accepted,
            )
        assert call.tool_name == "publish_artifact"
        return ToolResult(
            tool_use_id=call.tool_use_id,
            tool_name=call.tool_name,
            content=json.dumps({"status": "published"}),
            artifacts=[
                {
                    "id": "art-late-publish",
                    "kind": "artifact_ref",
                    "name": "late.html",
                    "mime": "text/html",
                    "size": 8,
                    "sha256": "a" * 64,
                    "session_id": "session-goal-terminal-boundary",
                    "session_key": "agent:main:webchat:goal-terminal-boundary",
                    "source": "publish_artifact",
                    "created_at": "2026-08-08T00:00:00Z",
                    "download_url": "/api/v1/artifacts/art-late-publish",
                }
            ],
        )

    agent = Agent(
        provider=provider,
        config=AgentConfig(
            max_provider_retries=0,
            max_turn_llm_calls=max_turn_llm_calls,
            max_turn_tool_errors=1 if terminal_accepted else 0,
            reasoning_stream_char_cap=5,
            thinking=summary_mode in {"reasoning_stream", "thinking_error"},
        ),
        tool_definitions=[
            ToolDefinition(
                name="update_goal",
                description="Complete the Goal",
                input_schema=ToolInputSchema(),
            ),
            ToolDefinition(
                name="publish_artifact",
                description="Publish a generated artifact",
                input_schema=ToolInputSchema(),
            ),
        ],
        tool_handler=handle_tool,
        tool_context=ToolContext(
            is_owner=True,
            caller_kind=CallerKind.WEB,
            workspace_dir=str(tmp_path),
            session_key="agent:main:webchat:goal-terminal-boundary",
            task_id="task-goal-terminal-boundary",
            collaboration_mode="default",
            goal_context={
                "schemaVersion": 1,
                "sessionId": "session-goal-terminal-boundary",
                "epoch": 0,
                "goalId": "goal-terminal-boundary",
                "objectiveRevision": 1,
                "objectiveSnapshot": "Complete before publishing anything else.",
                "taskId": "task-goal-terminal-boundary",
                "continuationSeq": 0,
                "automatic": False,
                "progress": None,
            },
        ),
    )

    events = [event async for event in agent.run_turn("finish the Goal")]

    summary_has_headroom = not terminal_accepted or max_turn_llm_calls == 0
    assert provider.calls == (2 if summary_has_headroom else 1)
    assert not any(isinstance(event, ErrorEvent) for event in events)
    if terminal_accepted:
        assert executed == ["update_goal"]
        if summary_has_headroom:
            assert provider.tools_by_call[1] == []
        assert not any(isinstance(event, ArtifactEvent) for event in events)
        skipped = next(
            event
            for event in events
            if isinstance(event, ToolResultEvent)
            and event.tool_name == "publish_artifact"
        )
        assert json.loads(str(skipped.result))["status"] == "not_executed"
        assert skipped.is_error is True
        done = next(event for event in events if isinstance(event, DoneEvent))
        assert done.text == (
            "Final Goal summary."
            if summary_has_headroom and summary_mode == "normal"
            else "The Goal is complete."
        )
    else:
        assert executed == ["update_goal", "publish_artifact"]
        assert set(provider.tools_by_call[1]) == {
            "update_goal",
            "publish_artifact",
        }
        assert any(isinstance(event, ArtifactEvent) for event in events)


@pytest.mark.asyncio
async def test_sessions_yield_is_a_dispatch_boundary_before_goal_terminal(
    tmp_path,
) -> None:
    provider = _GoalYieldThenTerminalProvider()
    executed: list[str] = []

    async def handle_tool(call: ToolCall) -> ToolResult:
        executed.append(call.tool_name)
        if call.tool_name == "sessions_yield":
            return ToolResult(
                tool_use_id=call.tool_use_id,
                tool_name=call.tool_name,
                content=json.dumps({"status": "yielded"}),
            )
        return ToolResult(
            tool_use_id=call.tool_use_id,
            tool_name=call.tool_name,
            content=json.dumps({"status": "accepted", "goal": {"status": "complete"}}),
        )

    agent = Agent(
        provider=provider,
        config=AgentConfig(max_provider_retries=0),
        tool_definitions=[
            ToolDefinition(
                name="sessions_yield",
                description="Yield the current turn",
                input_schema=ToolInputSchema(),
            ),
            ToolDefinition(
                name="update_goal",
                description="Complete the Goal",
                input_schema=ToolInputSchema(),
            ),
        ],
        tool_handler=handle_tool,
        tool_context=ToolContext(
            is_owner=True,
            caller_kind=CallerKind.WEB,
            workspace_dir=str(tmp_path),
            session_key="agent:main:webchat:goal-yield-boundary",
            task_id="task-goal-yield-boundary",
            collaboration_mode="default",
            goal_context={
                "schemaVersion": 1,
                "sessionId": "session-goal-yield-boundary",
                "epoch": 0,
                "goalId": "goal-yield-boundary",
                "objectiveRevision": 1,
                "objectiveSnapshot": "Yield before any later terminal decision.",
                "taskId": "task-goal-yield-boundary",
                "continuationSeq": 0,
                "automatic": False,
                "progress": None,
            },
        ),
    )

    events = [event async for event in agent.run_turn("yield now")]

    assert provider.calls == 1
    assert executed == ["sessions_yield"]
    skipped = next(
        event
        for event in events
        if isinstance(event, ToolResultEvent) and event.tool_name == "update_goal"
    )
    assert json.loads(str(skipped.result))["status"] == "not_executed"
    assert not any(isinstance(event, ErrorEvent) for event in events)


@pytest.mark.asyncio
async def test_turn_runner_auto_publishes_deliverable_file_when_model_omits_publish(
    tmp_path,
) -> None:
    storage = SessionStorage(":memory:")
    await storage.connect()
    manager = SessionManager(storage)
    session_key = "agent:main:webchat:artifact-omitted"
    session = await manager.create(session_key)
    runner = TurnRunner(
        provider_selector=_ProviderSelector(_OmittedPublishProvider()),
        tool_registry=_write_file_registry(),
        session_manager=manager,
        config=GatewayConfig(
            attachments=AttachmentsConfig(media_root=str(tmp_path / "media")),
            squilla_router=SquillaRouterConfig(enabled=False),
        ),
    )
    tool_context = ToolContext(
        is_owner=True,
        caller_kind=CallerKind.WEB,
        workspace_dir=str(tmp_path / "workspace"),
        allowed_tools={"write_file"},
        elevated="full",
    )

    try:
        events = [
            event
            async for event in runner.run(
                "make an html page",
                session_key,
                tool_context=tool_context,
                history_has_persisted_user=False,
                no_memory_capture=True,
            )
        ]

        artifact_events = [event for event in events if isinstance(event, ArtifactEvent)]
        assert len(artifact_events) == 1
        assert artifact_events[0].name == "manual-big-write.html"
        assert artifact_events[0].mime == "text/html"
        assert artifact_events[0].session_id == session.session_id
        assert artifact_events[0].download_url == (
            f"/api/v1/artifacts/{artifact_events[0].id}"
        )

        transcript = await manager.get_transcript(session_key)
        assistant = [entry for entry in transcript if entry.role == "assistant"][-1]
        payload = json.loads(assistant.content)
        assert payload["text"] == "Created manual-big-write.html for you."
        assert payload["artifacts"][0]["name"] == "manual-big-write.html"
        assert payload["artifacts"][0]["source"] == "auto_publish_omitted"
    finally:
        await storage.close()


@pytest.mark.asyncio
async def test_repeated_cancel_persists_completed_artifact_and_interrupted_transcript(
    tmp_path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    storage = SessionStorage(":memory:")
    await storage.connect()
    manager = SessionManager(storage)
    session_key = "agent:main:webchat:artifact-repeated-cancel"
    session = await manager.create(session_key)
    runner = TurnRunner(
        provider_selector=_ProviderSelector(_OmittedPublishProvider()),
        tool_registry=_write_file_registry(),
        session_manager=manager,
        config=GatewayConfig(
            attachments=AttachmentsConfig(media_root=str(tmp_path / "media")),
            squilla_router=SquillaRouterConfig(enabled=False),
        ),
    )
    tool_context = ToolContext(
        is_owner=True,
        caller_kind=CallerKind.WEB,
        workspace_dir=str(tmp_path / "workspace"),
        allowed_tools={"write_file"},
        elevated="full",
    )

    publish_started = threading.Event()
    release_publish = threading.Event()
    publish_finished = threading.Event()
    real_publish = runner._stream_consumer_stage._done_handler.run_publish

    def gated_publish(inp: Any, accumulated_text: str) -> Any:
        publish_started.set()
        assert release_publish.wait(timeout=5.0), "publish was never released"
        result = real_publish(inp, accumulated_text)
        publish_finished.set()
        return result

    monkeypatch.setattr(
        runner._stream_consumer_stage._done_handler,
        "run_publish",
        gated_publish,
    )

    assistant_append_started = asyncio.Event()
    release_assistant_append = asyncio.Event()
    real_append_message = manager.append_message

    async def gated_append_message(session_key: str, **kwargs: Any) -> Any:
        if kwargs.get("role") == "assistant":
            assistant_append_started.set()
            await release_assistant_append.wait()
        return await real_append_message(session_key, **kwargs)

    monkeypatch.setattr(manager, "append_message", gated_append_message)

    async def consume() -> None:
        async for _event in runner.run(
            "make an html page",
            session_key,
            tool_context=tool_context,
            history_has_persisted_user=False,
            no_memory_capture=True,
        ):
            pass

    task = asyncio.create_task(consume())
    try:
        assert await asyncio.to_thread(publish_started.wait, 5.0)
        task.cancel("cancel-during-publish")
        release_publish.set()

        await asyncio.wait_for(assistant_append_started.wait(), timeout=5.0)
        assert publish_finished.is_set()
        task.cancel("cancel-during-transcript")
        await asyncio.sleep(0)

        release_assistant_append.set()
        with pytest.raises(asyncio.CancelledError) as exc_info:
            await task
        assert exc_info.value.args == ("cancel-during-publish",)

        transcript = await manager.get_transcript(session_key)
        assistants = [entry for entry in transcript if entry.role == "assistant"]
        assert len(assistants) == 1
        payload = json.loads(assistants[0].content)
        assert "The generated file was delivered" in payload["text"]
        assert payload["artifacts"][0]["name"] == "manual-big-write.html"

        store = ArtifactStore(str(tmp_path / "media"))
        _, artifact_path = store.resolve_for_download(
            payload["artifacts"][0]["id"],
            session_id=session.session_id,
        )
        assert artifact_path.read_text(encoding="utf-8") == (
            "<!doctype html><title>Manual</title>"
        )
    finally:
        release_publish.set()
        release_assistant_append.set()
        if not task.done():
            task.cancel()
            with contextlib.suppress(asyncio.CancelledError):
                await task
        await storage.close()


def test_auto_publish_validates_and_publishes_pptx_from_same_bytes(
    tmp_path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    import openstarry_code.engine.artifact_delivery as artifact_delivery

    workspace = tmp_path / "workspace"
    workspace.mkdir()
    target = workspace / "brief.pptx"
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[1])
    presentation.save(target)
    expected = target.read_bytes()
    replacement = Presentation()
    replacement.slides.add_slide(replacement.slide_layouts[5])
    replacement_output = BytesIO()
    replacement.save(replacement_output)
    replacement_payload = replacement_output.getvalue()
    validate = artifact_delivery.validate_artifact_for_delivery

    def validate_then_replace(*args: object, **kwargs: object) -> object:
        report = validate(*args, **kwargs)
        target.write_bytes(replacement_payload)
        return report

    monkeypatch.setattr(
        artifact_delivery,
        "validate_artifact_for_delivery",
        validate_then_replace,
    )

    ctx = ToolContext(
        is_owner=True,
        caller_kind=CallerKind.WEB,
        workspace_dir=str(workspace),
        artifact_media_root=str(tmp_path / "media"),
        artifact_session_id="session-valid-pptx",
        session_key="agent:main:webchat:valid-pptx",
    )
    ctx.workspace_file_writes.append(
        {
            "created": True,
            "path": str(target),
            "relative_path": target.name,
            "name": target.name,
        }
    )

    result = auto_publish_omitted_workspace_artifacts(
        ctx,
        final_text="Created brief.pptx for you.",
    )

    assert result.failure_summaries == []
    assert len(result.artifacts) == 1
    artifact = result.artifacts[0]
    assert artifact["name"] == "brief.pptx"
    store = ArtifactStore(str(tmp_path / "media"))
    _, material_path = store.resolve_for_download(
        str(artifact["id"]),
        session_id="session-valid-pptx",
    )
    material = material_path.read_bytes()
    assert target.read_bytes() == replacement_payload
    assert material == expected
    Presentation(BytesIO(material))


def test_auto_publish_validates_invalid_pptx_before_persisted_dedupe(tmp_path) -> None:
    workspace = tmp_path / "workspace"
    workspace.mkdir()
    target = workspace / "brief.pptx"
    invalid_payload = b"not an OOXML package"
    target.write_bytes(invalid_payload)
    media_root = tmp_path / "media"
    store = ArtifactStore(media_root)
    historical = store.publish_bytes(
        invalid_payload,
        session_id="session-invalid-pptx",
        session_key="agent:main:webchat:invalid-pptx",
        name=target.name,
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        source="legacy",
    )
    ctx = ToolContext(
        is_owner=True,
        caller_kind=CallerKind.WEB,
        workspace_dir=str(workspace),
        artifact_media_root=str(media_root),
        artifact_session_id="session-invalid-pptx",
        session_key="agent:main:webchat:invalid-pptx",
    )
    ctx.workspace_file_writes.append(
        {
            "created": True,
            "path": str(target),
            "relative_path": target.name,
            "name": target.name,
        }
    )

    result = auto_publish_omitted_workspace_artifacts(
        ctx,
        final_text="Created brief.pptx for you.",
    )

    assert result.artifacts == []
    assert len(result.failure_summaries) == 1
    assert ctx.published_artifacts == []
    assert store.path_for(historical).read_bytes() == invalid_payload


def test_auto_publish_known_valid_pptx_reports_exact_resolved_target(tmp_path) -> None:
    workspace = tmp_path / "workspace"
    reports = workspace / "reports"
    reports.mkdir(parents=True)
    target = reports / "brief.pptx"
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[1])
    presentation.save(target)
    target_sha256 = hashlib.sha256(target.read_bytes()).hexdigest()
    ctx = ToolContext(
        is_owner=True,
        caller_kind=CallerKind.WEB,
        workspace_dir=str(workspace),
        artifact_media_root=str(tmp_path / "media"),
        artifact_session_id="session-known-pptx",
        session_key="agent:main:webchat:known-pptx",
    )
    ctx.published_artifacts.append(
        {
            "id": "already-present",
            "sha256": target_sha256,
            "name": target.name,
        }
    )
    ctx.workspace_file_writes.append(
        {
            "created": True,
            "path": str(target),
            "relative_path": "reports/brief.pptx",
            "name": target.name,
        }
    )

    result = auto_publish_omitted_workspace_artifacts(
        ctx,
        final_text="Created reports/brief.pptx for you.",
    )

    assert result.artifacts == []
    assert result.failure_summaries == []
    assert set(result.resolved_target_keys) == {
        "path:" + os.path.normcase(os.path.normpath(str(target.resolve()))),
        "name:brief.pptx",
    }


def test_auto_publish_nested_target_does_not_report_basename_as_resolved(tmp_path) -> None:
    workspace = tmp_path / "workspace"
    reports = workspace / "reports"
    reports.mkdir(parents=True)
    target = reports / "deck.html"
    target.write_text("<title>Deck</title>", encoding="utf-8")
    ctx = ToolContext(
        is_owner=True,
        caller_kind=CallerKind.WEB,
        workspace_dir=str(workspace),
        artifact_media_root=str(tmp_path / "media"),
        artifact_session_id="session-nested-target",
        session_key="agent:main:webchat:nested-target",
    )
    ctx.workspace_file_writes.append(
        {
            "created": True,
            "path": str(target),
            "relative_path": "reports/deck.html",
            "name": target.name,
        }
    )

    result = auto_publish_omitted_workspace_artifacts(
        ctx,
        final_text="Created reports/deck.html for you.",
    )

    assert len(result.artifacts) == 1
    assert set(result.resolved_target_keys) == {
        "path:" + os.path.normcase(os.path.normpath(str(target.resolve()))),
        "name:deck.html",
    }
    assert "path:" + os.path.normcase(os.path.normpath("deck.html")) not in (
        result.resolved_target_keys
    )


def test_artifact_delivery_target_key_preserves_whitespace_and_tolerates_nul(
    tmp_path,
) -> None:
    workspace = tmp_path / "workspace"

    plain = artifact_delivery_publish_target_key(
        "deck.pptx",
        workspace_dir=workspace,
    )
    leading_space = artifact_delivery_publish_target_key(
        " deck.pptx",
        workspace_dir=workspace,
    )

    assert plain is not None
    assert leading_space is not None
    assert plain != leading_space
    assert artifact_delivery_publish_target_key(
        "bad\x00deck.pptx",
        workspace_dir=workspace,
    ) is None


@pytest.mark.parametrize(
    "target_name",
    [
        pytest.param(
            "deck:unsafe.html",
            marks=pytest.mark.skipif(
                os.name == "nt",
                reason="colon is not a legal Windows workspace filename",
            ),
        ),
        ("x" * 156) + ".html",
    ],
)
def test_auto_publish_dedupes_current_artifact_by_store_safe_name(
    tmp_path_factory: pytest.TempPathFactory,
    target_name: str,
) -> None:
    root = tmp_path_factory.mktemp("a")
    workspace = root / "w"
    workspace.mkdir()
    target = workspace / target_name
    payload = b"<title>Already published</title>"
    target.write_bytes(payload)
    media_root = root / "m"
    store = ArtifactStore(media_root)
    existing = store.publish_bytes(
        payload,
        session_id="session-safe-name",
        session_key="agent:main:webchat:safe-name",
        name=target.name,
        mime="text/html",
        source="publish_artifact",
    )
    ctx = ToolContext(
        is_owner=True,
        caller_kind=CallerKind.WEB,
        workspace_dir=str(workspace),
        artifact_media_root=str(media_root),
        artifact_session_id="session-safe-name",
        session_key="agent:main:webchat:safe-name",
        published_artifacts=[artifact_payload(existing)],
    )
    ctx.workspace_file_writes.append(
        {
            "created": True,
            "path": str(target),
            "relative_path": target.name,
            "name": target.name,
        }
    )

    result = auto_publish_omitted_workspace_artifacts(
        ctx,
        final_text=f"Created {target.name} for you.",
    )

    assert result.artifacts == []
    assert result.failure_summaries == []
    assert len(ctx.published_artifacts) == 1


def test_auto_publish_oversized_pptx_preflights_before_read_or_validation(
    tmp_path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    import openstarry_code.engine.artifact_delivery as artifact_delivery

    workspace = tmp_path / "workspace"
    workspace.mkdir()
    target = workspace / "large.pptx"
    target.write_bytes(b"oversized")
    ctx = ToolContext(
        is_owner=True,
        caller_kind=CallerKind.WEB,
        workspace_dir=str(workspace),
        artifact_media_root=str(tmp_path / "media"),
        artifact_session_id="session-large-pptx",
        session_key="agent:main:webchat:large-pptx",
        artifact_max_bytes=4,
    )
    ctx.workspace_file_writes.append(
        {
            "created": True,
            "path": str(target),
            "relative_path": target.name,
            "name": target.name,
        }
    )

    def unexpected_call(*args: object, **kwargs: object) -> None:
        pytest.fail("oversized PPTX must be rejected before reading or validation")

    monkeypatch.setattr(artifact_delivery.Path, "read_bytes", unexpected_call)
    monkeypatch.setattr(
        artifact_delivery,
        "validate_artifact_for_delivery",
        unexpected_call,
    )

    result = auto_publish_omitted_workspace_artifacts(
        ctx,
        final_text="Created large.pptx for you.",
    )

    assert result.artifacts == []
    assert len(result.failure_summaries) == 1
    assert "per-file budget" in result.failure_summaries[0]
    assert ctx.published_artifacts == []


@pytest.mark.asyncio
async def test_turn_runner_rejects_invalid_omitted_pptx_and_marks_delivery_failure(
    tmp_path,
) -> None:
    storage = SessionStorage(":memory:")
    await storage.connect()
    manager = SessionManager(storage)
    session_key = "agent:main:webchat:artifact-invalid-pptx"
    await manager.create(session_key)
    runner = TurnRunner(
        provider_selector=_ProviderSelector(_OmittedInvalidPptxProvider()),
        tool_registry=_write_file_registry(),
        session_manager=manager,
        config=GatewayConfig(
            attachments=AttachmentsConfig(media_root=str(tmp_path / "media")),
            squilla_router=SquillaRouterConfig(enabled=False),
        ),
    )
    tool_context = ToolContext(
        is_owner=True,
        caller_kind=CallerKind.WEB,
        workspace_dir=str(tmp_path / "workspace"),
        allowed_tools={"write_file"},
        elevated="full",
    )

    try:
        events = [
            event
            async for event in runner.run(
                "make a presentation",
                session_key,
                tool_context=tool_context,
                history_has_persisted_user=False,
                no_memory_capture=True,
            )
        ]

        assert [event for event in events if isinstance(event, ArtifactEvent)] == []
        done = next(event for event in events if isinstance(event, DoneEvent))
        assert "File delivery failed:" in done.text
        assert "correct or regenerate it" in done.text

        transcript = await manager.get_transcript(session_key)
        assistant = [entry for entry in transcript if entry.role == "assistant"][-1]
        assert "File delivery failed:" in assistant.content
        assert "artifacts" not in assistant.content
    finally:
        await storage.close()


@pytest.mark.asyncio
async def test_turn_runner_does_not_auto_publish_edited_config_json(tmp_path) -> None:
    storage = SessionStorage(":memory:")
    await storage.connect()
    manager = SessionManager(storage)
    session_key = "agent:main:webchat:artifact-edit-config"
    await manager.create(session_key)
    workspace = tmp_path / "workspace"
    workspace.mkdir()
    (workspace / "config.json").write_text("{\"enabled\": false}\n", encoding="utf-8")
    runner = TurnRunner(
        provider_selector=_ProviderSelector(_EditedConfigProvider()),
        tool_registry=_edit_file_registry(),
        session_manager=manager,
        config=GatewayConfig(
            attachments=AttachmentsConfig(media_root=str(tmp_path / "media")),
            squilla_router=SquillaRouterConfig(enabled=False),
        ),
    )
    tool_context = ToolContext(
        is_owner=True,
        caller_kind=CallerKind.WEB,
        workspace_dir=str(workspace),
        allowed_tools={"edit_file"},
        elevated="full",
    )

    try:
        events = [
            event
            async for event in runner.run(
                "update config",
                session_key,
                tool_context=tool_context,
                history_has_persisted_user=False,
                no_memory_capture=True,
            )
        ]

        artifact_events = [event for event in events if isinstance(event, ArtifactEvent)]
        assert artifact_events == []

        transcript = await manager.get_transcript(session_key)
        assistant = [entry for entry in transcript if entry.role == "assistant"][-1]
        assert assistant.content == "Updated config.json."
    finally:
        await storage.close()


@pytest.mark.asyncio
async def test_turn_runner_auto_publishes_deliverable_file_created_by_apply_patch(
    tmp_path,
) -> None:
    storage = SessionStorage(":memory:")
    await storage.connect()
    manager = SessionManager(storage)
    session_key = "agent:main:webchat:artifact-patch-omitted"
    await manager.create(session_key)
    runner = TurnRunner(
        provider_selector=_ProviderSelector(_OmittedPatchPublishProvider()),
        tool_registry=_apply_patch_registry(),
        session_manager=manager,
        config=GatewayConfig(
            attachments=AttachmentsConfig(media_root=str(tmp_path / "media")),
            squilla_router=SquillaRouterConfig(enabled=False),
        ),
    )
    tool_context = ToolContext(
        is_owner=True,
        caller_kind=CallerKind.WEB,
        workspace_dir=str(tmp_path / "workspace"),
        allowed_tools={"apply_patch"},
        elevated="full",
    )

    try:
        events = [
            event
            async for event in runner.run(
                "make an html page with a patch",
                session_key,
                tool_context=tool_context,
                history_has_persisted_user=False,
                no_memory_capture=True,
            )
        ]

        artifact_events = [event for event in events if isinstance(event, ArtifactEvent)]
        assert [event.name for event in artifact_events] == ["patched.html"]

        transcript = await manager.get_transcript(session_key)
        assistant = [entry for entry in transcript if entry.role == "assistant"][-1]
        payload = json.loads(assistant.content)
        assert payload["artifacts"][0]["name"] == "patched.html"
    finally:
        await storage.close()


@pytest.mark.asyncio
async def test_turn_runner_marks_partial_omitted_artifact_delivery_failure(
    tmp_path,
) -> None:
    storage = SessionStorage(":memory:")
    await storage.connect()
    manager = SessionManager(storage)
    session_key = "agent:main:webchat:artifact-partial-failure"
    await manager.create(session_key)
    runner = TurnRunner(
        provider_selector=_ProviderSelector(_MixedSizeOmittedPublishProvider()),
        tool_registry=_write_file_registry(),
        session_manager=manager,
        config=GatewayConfig(
            attachments=AttachmentsConfig(
                media_root=str(tmp_path / "media"),
                artifact_max_bytes=40,
            ),
            squilla_router=SquillaRouterConfig(enabled=False),
        ),
    )
    tool_context = ToolContext(
        is_owner=True,
        caller_kind=CallerKind.WEB,
        workspace_dir=str(tmp_path / "workspace"),
        allowed_tools={"write_file"},
        elevated="full",
    )

    try:
        events = [
            event
            async for event in runner.run(
                "make two html pages",
                session_key,
                tool_context=tool_context,
                history_has_persisted_user=False,
                no_memory_capture=True,
            )
        ]

        artifact_events = [event for event in events if isinstance(event, ArtifactEvent)]
        assert [event.name for event in artifact_events] == ["small.html"]
        text_deltas = [event.text for event in events if isinstance(event, TextDeltaEvent)]
        done = next(event for event in events if isinstance(event, DoneEvent))
        assert any("File delivery failed:" in text for text in text_deltas)
        assert "File delivery failed:" in done.text
        assert "some generated files were attached" in done.text
        assert "correct or regenerate it" in done.text
        assert "no downloadable file was attached" not in done.text

        transcript = await manager.get_transcript(session_key)
        assistant = [entry for entry in transcript if entry.role == "assistant"][-1]
        payload = json.loads(assistant.content)
        assert payload["artifacts"][0]["name"] == "small.html"
        assert "File delivery failed:" in payload["text"]
        assert "some generated files were attached" in payload["text"]
        assert "no downloadable file was attached" not in payload["text"]
    finally:
        await storage.close()


@pytest.mark.asyncio
async def test_turn_runner_auto_publishes_same_content_deliverables_by_name(
    tmp_path,
) -> None:
    storage = SessionStorage(":memory:")
    await storage.connect()
    manager = SessionManager(storage)
    session_key = "agent:main:webchat:artifact-same-content"
    await manager.create(session_key)
    runner = TurnRunner(
        provider_selector=_ProviderSelector(_SameContentOmittedPublishProvider()),
        tool_registry=_write_file_registry(),
        session_manager=manager,
        config=GatewayConfig(
            attachments=AttachmentsConfig(media_root=str(tmp_path / "media")),
            squilla_router=SquillaRouterConfig(enabled=False),
        ),
    )
    tool_context = ToolContext(
        is_owner=True,
        caller_kind=CallerKind.WEB,
        workspace_dir=str(tmp_path / "workspace"),
        allowed_tools={"write_file"},
        elevated="full",
    )

    try:
        events = [
            event
            async for event in runner.run(
                "make matching html files",
                session_key,
                tool_context=tool_context,
                history_has_persisted_user=False,
                no_memory_capture=True,
            )
        ]

        artifact_events = [event for event in events if isinstance(event, ArtifactEvent)]
        assert [event.name for event in artifact_events] == ["first.html", "second.html"]
        assert artifact_events[0].sha256 == artifact_events[1].sha256

        transcript = await manager.get_transcript(session_key)
        assistant = [entry for entry in transcript if entry.role == "assistant"][-1]
        payload = json.loads(assistant.content)
        assert [artifact["name"] for artifact in payload["artifacts"]] == [
            "first.html",
            "second.html",
        ]
    finally:
        await storage.close()


@pytest.mark.asyncio
async def test_turn_runner_auto_publishes_omitted_deliverable_after_existing_artifact(
    tmp_path,
) -> None:
    storage = SessionStorage(":memory:")
    await storage.connect()
    manager = SessionManager(storage)
    session_key = "agent:main:webchat:artifact-partial-omitted"
    session = await manager.create(session_key)
    runner = TurnRunner(
        provider_selector=_ProviderSelector(_PartialOmittedPublishProvider()),
        tool_registry=_registry_with_write_file(),
        session_manager=manager,
        config=GatewayConfig(
            attachments=AttachmentsConfig(media_root=str(tmp_path / "media")),
            squilla_router=SquillaRouterConfig(enabled=False),
        ),
    )
    tool_context = ToolContext(
        is_owner=True,
        caller_kind=CallerKind.WEB,
        workspace_dir=str(tmp_path / "workspace"),
        allowed_tools={"make_file", "write_file"},
        elevated="full",
    )

    try:
        events = [
            event
            async for event in runner.run(
                "make two files",
                session_key,
                tool_context=tool_context,
                history_has_persisted_user=False,
                no_memory_capture=True,
            )
        ]

        artifact_events = [event for event in events if isinstance(event, ArtifactEvent)]
        assert [event.name for event in artifact_events] == ["runtime.txt", "second.html"]
        assert artifact_events[0].id == "art-runtime"
        assert artifact_events[1].session_id == session.session_id

        transcript = await manager.get_transcript(session_key)
        assistant = [entry for entry in transcript if entry.role == "assistant"][-1]
        payload = json.loads(assistant.content)
        assert [artifact["name"] for artifact in payload["artifacts"]] == [
            "runtime.txt",
            "second.html",
        ]
    finally:
        await storage.close()


@pytest.mark.asyncio
async def test_turn_runner_does_not_auto_publish_memory_json_write(tmp_path) -> None:
    storage = SessionStorage(":memory:")
    await storage.connect()
    manager = SessionManager(storage)
    session_key = "agent:main:webchat:artifact-memory-json"
    await manager.create(session_key)
    runner = TurnRunner(
        provider_selector=_ProviderSelector(_MemoryJsonWriteProvider()),
        tool_registry=_write_file_registry(),
        session_manager=manager,
        config=GatewayConfig(
            attachments=AttachmentsConfig(media_root=str(tmp_path / "media")),
            squilla_router=SquillaRouterConfig(enabled=False),
        ),
    )
    tool_context = ToolContext(
        is_owner=True,
        caller_kind=CallerKind.WEB,
        workspace_dir=str(tmp_path / "workspace"),
        allowed_tools={"write_file"},
        elevated="full",
    )

    try:
        events = [
            event
            async for event in runner.run(
                "update internal memory",
                session_key,
                tool_context=tool_context,
                history_has_persisted_user=False,
                no_memory_capture=True,
            )
        ]

        artifact_events = [event for event in events if isinstance(event, ArtifactEvent)]
        assert artifact_events == []

        transcript = await manager.get_transcript(session_key)
        assistant = [entry for entry in transcript if entry.role == "assistant"][-1]
        assert assistant.content == "Updated memory/cache.json."
    finally:
        await storage.close()


@pytest.mark.asyncio
async def test_turn_runner_marks_failed_omitted_artifact_delivery_in_final_text(
    tmp_path,
) -> None:
    storage = SessionStorage(":memory:")
    await storage.connect()
    manager = SessionManager(storage)
    session_key = "agent:main:webchat:artifact-omitted-too-large"
    await manager.create(session_key)
    runner = TurnRunner(
        provider_selector=_ProviderSelector(_OmittedPublishProvider()),
        tool_registry=_write_file_registry(),
        session_manager=manager,
        config=GatewayConfig(
            attachments=AttachmentsConfig(
                media_root=str(tmp_path / "media"),
                artifact_max_bytes=1,
            ),
            squilla_router=SquillaRouterConfig(enabled=False),
        ),
    )
    tool_context = ToolContext(
        is_owner=True,
        caller_kind=CallerKind.WEB,
        workspace_dir=str(tmp_path / "workspace"),
        allowed_tools={"write_file"},
        elevated="full",
    )

    try:
        events = [
            event
            async for event in runner.run(
                "make an html page",
                session_key,
                tool_context=tool_context,
                history_has_persisted_user=False,
                no_memory_capture=True,
            )
        ]

        artifact_events = [event for event in events if isinstance(event, ArtifactEvent)]
        assert artifact_events == []
        text_deltas = [event.text for event in events if isinstance(event, TextDeltaEvent)]
        done = next(event for event in events if isinstance(event, DoneEvent))
        assert any("File delivery failed:" in text for text in text_deltas)
        assert "File delivery failed:" in done.text

        transcript = await manager.get_transcript(session_key)
        assistant = [entry for entry in transcript if entry.role == "assistant"][-1]
        assert "File delivery failed:" in assistant.content
        assert "artifacts" not in assistant.content
    finally:
        await storage.close()


@pytest.mark.asyncio
async def test_turn_runner_marks_failed_artifact_delivery_in_final_text(tmp_path) -> None:
    storage = SessionStorage(":memory:")
    await storage.connect()
    manager = SessionManager(storage)
    session_key = "agent:main:webchat:artifact-failed"
    await manager.create(session_key)
    runner = TurnRunner(
        provider_selector=_ProviderSelector(_FailedPublishProvider()),
        tool_registry=_failed_publish_registry(),
        session_manager=manager,
        config=GatewayConfig(
            attachments=AttachmentsConfig(media_root=str(tmp_path / "media")),
            squilla_router=SquillaRouterConfig(enabled=False),
        ),
    )
    tool_context = ToolContext(
        is_owner=True,
        caller_kind=CallerKind.WEB,
        workspace_dir=str(tmp_path),
    )

    try:
        events = [
            event
            async for event in runner.run(
                "make report",
                session_key,
                tool_context=tool_context,
                history_has_persisted_user=False,
                no_memory_capture=True,
            )
        ]

        text_deltas = [event.text for event in events if isinstance(event, TextDeltaEvent)]
        done = next(event for event in events if isinstance(event, DoneEvent))
        assert any("File delivery failed:" in text for text in text_deltas)
        assert "File delivery failed:" in done.text
        assert "Ask me to resend the file after I correct or regenerate it." in done.text
        assert "publish_artifact" not in done.text
        assert "active workspace" not in done.text
        assert "missing-report.pptx" not in done.text

        transcript = await manager.get_transcript(session_key)
        assistant = [entry for entry in transcript if entry.role == "assistant"][-1]
        assert "Report file is ready for download." in assistant.content
        assert "File delivery failed:" in assistant.content
        assert "artifacts" not in assistant.content
    finally:
        await storage.close()


@pytest.mark.asyncio
async def test_turn_runner_clears_delivery_failure_after_same_target_retry_succeeds(
    tmp_path,
) -> None:
    storage = SessionStorage(":memory:")
    await storage.connect()
    manager = SessionManager(storage)
    session_key = "agent:main:webchat:artifact-retry-success"
    await manager.create(session_key)
    provider = _RetryPublishProvider()
    runner = TurnRunner(
        provider_selector=_ProviderSelector(provider),
        tool_registry=_retry_publish_registry(),
        session_manager=manager,
        config=GatewayConfig(
            attachments=AttachmentsConfig(media_root=str(tmp_path / "media")),
            squilla_router=SquillaRouterConfig(enabled=False),
        ),
    )
    tool_context = ToolContext(
        is_owner=True,
        caller_kind=CallerKind.WEB,
        workspace_dir=str(tmp_path),
    )

    try:
        events = [
            event
            async for event in runner.run(
                "make report",
                session_key,
                tool_context=tool_context,
                history_has_persisted_user=False,
                no_memory_capture=True,
            )
        ]

        done = next(event for event in events if isinstance(event, DoneEvent))
        artifacts = [event for event in events if isinstance(event, ArtifactEvent)]
        assert provider.calls == 2
        assert [artifact.id for artifact in artifacts] == ["art-retried"]
        assert "File delivery failed:" not in done.text

        transcript = await manager.get_transcript(session_key)
        assistant = [entry for entry in transcript if entry.role == "assistant"][-1]
        assert "File delivery failed:" not in assistant.content
        assert "art-retried" in assistant.content
    finally:
        await storage.close()


@pytest.mark.asyncio
async def test_turn_runner_marks_failed_create_pptx_delivery_in_final_text(tmp_path) -> None:
    storage = SessionStorage(":memory:")
    await storage.connect()
    manager = SessionManager(storage)
    session_key = "agent:main:webchat:create-pptx-failed"
    await manager.create(session_key)
    runner = TurnRunner(
        provider_selector=_ProviderSelector(_FailedCreatePptxProvider()),
        tool_registry=_failed_create_pptx_registry(),
        session_manager=manager,
        config=GatewayConfig(
            attachments=AttachmentsConfig(media_root=str(tmp_path / "media")),
            squilla_router=SquillaRouterConfig(enabled=False),
        ),
    )
    tool_context = ToolContext(
        is_owner=True,
        caller_kind=CallerKind.WEB,
        workspace_dir=str(tmp_path),
    )

    try:
        events = [
            event
            async for event in runner.run(
                "make report",
                session_key,
                tool_context=tool_context,
                history_has_persisted_user=False,
                no_memory_capture=True,
            )
        ]
        done = next(event for event in events if isinstance(event, DoneEvent))
        assert [event for event in events if isinstance(event, ArtifactEvent)] == []
        assert "File delivery failed:" in done.text
        assert "correct or regenerate it" in done.text
    finally:
        await storage.close()


class _GoalArtifactTopologySelector:
    """Minimal per-turn selector used to exercise the real TurnRunner topology."""

    def __init__(self, provider: Any) -> None:
        self.provider = provider
        self.clone_calls = 0
        self.current_config = SimpleNamespace(provider="test", model="deepseek-v4-pro")

    def clone(self) -> _GoalArtifactTopologySelector:
        self.clone_calls += 1
        return self

    @property
    def active_provider_id(self) -> str:
        return str(self.current_config.provider)

    def override_model(self, model: str) -> None:
        self.current_config = SimpleNamespace(provider="test", model=model)
        self.provider.model = model

    def remaining_chain(self) -> list[SimpleNamespace]:
        return [self.current_config]

    def resolve(self) -> Any:
        return self.provider


class _GoalArtifactPrimaryFailureProvider:
    provider_name = "test-primary"

    def __init__(self) -> None:
        self.calls = 0
        self.model = "test/primary"

    def chat(self, messages: list[Message], tools=None, config=None) -> AsyncIterator[Any]:
        self.calls += 1
        return self._stream()

    async def _stream(self) -> AsyncIterator[Any]:
        yield ProviderError(message="synthetic primary failure", code="503")

    async def list_models(self) -> list[ModelInfo]:
        return []


class _GoalArtifactFallbackSelector:
    """Fail before content once, then keep the fallback active for the whole loop."""

    def __init__(self, fallback: _GoalPostPublishLoopProvider) -> None:
        self.primary = _GoalArtifactPrimaryFailureProvider()
        self.fallback = fallback
        self.clone_calls = 0
        self.fallback_calls = 0
        self.current_config = SimpleNamespace(
            provider="test-primary",
            model="deepseek-v4-pro",
        )
        self._remaining = [
            self.current_config,
            SimpleNamespace(provider="test-fallback", model="deepseek-v4-flash"),
        ]

    def clone(self) -> _GoalArtifactFallbackSelector:
        self.clone_calls += 1
        return self

    @property
    def active_provider_id(self) -> str:
        return str(self.current_config.provider)

    def override_model(self, model: str) -> None:
        assert model == self.current_config.model

    def remaining_chain(self) -> list[SimpleNamespace]:
        return list(self._remaining)

    def resolve(self) -> Any:
        if self.current_config.model == "deepseek-v4-flash":
            return self.fallback
        return self.primary

    def next_fallback_after_failure(
        self,
        exc: Exception,
    ) -> _GoalPostPublishLoopProvider:
        del exc
        self.fallback_calls += 1
        self.current_config = self._remaining[1]
        self._remaining = self._remaining[1:]
        return self.fallback


class _GoalArtifactEnsembleProvider(_GoalPostPublishLoopProvider):
    """Attach one proposer and aggregator receipt to each normal decision."""

    provider_name = "ensemble"

    async def _stream(self, call_number: int) -> AsyncIterator[Any]:
        async for event in super()._stream(call_number):
            if isinstance(event, ProviderDone):
                event.model = "test/aggregator"
                event.input_tokens = 3
                event.output_tokens = 2
                event.model_usage_breakdown = [
                    {
                        "role": "proposer",
                        "label": "proposer_1",
                        "provider": "test",
                        "model": "test/proposer",
                        "input_tokens": 1,
                        "output_tokens": 1,
                    },
                    {
                        "role": "aggregator",
                        "label": "aggregator",
                        "provider": "test",
                        "model": "test/aggregator",
                        "input_tokens": 2,
                        "output_tokens": 1,
                    },
                ]
                event.ensemble_trace = {
                    "profile": "test",
                    "llm_request_count": 2,
                }
            yield event


async def _run_goal_artifact_topology(
    tmp_path,
    *,
    provider: _GoalPostPublishLoopProvider,
    selector: Any,
    config: GatewayConfig,
) -> tuple[list[str], list[str], list[Any]]:
    storage = SessionStorage(":memory:")
    await storage.connect()
    manager = SessionManager(storage)
    session_key = "agent:main:webchat:goal-artifact-topology"
    await manager.create(session_key)
    registry, control_calls, qa_calls = _goal_publish_loop_registry()
    runner = TurnRunner(
        provider_selector=selector,
        tool_registry=registry,
        session_manager=manager,
        config=config,
    )
    tool_context = ToolContext(
        is_owner=True,
        caller_kind=CallerKind.WEB,
        workspace_dir=str(tmp_path),
        session_key=session_key,
        task_id="task-goal-artifact-topology",
        collaboration_mode="default",
        goal_context={
            "schemaVersion": 1,
            "sessionId": "session-goal-artifact-topology",
            "epoch": 0,
            "goalId": "goal-artifact-topology",
            "objectiveRevision": 1,
            "objectiveSnapshot": "Publish and verify the report.",
            "taskId": "task-goal-artifact-topology",
            "continuationSeq": 0,
            "automatic": False,
            "progress": {
                "steps": [
                    {"step": "Publish the verified artifact", "status": "in_progress"}
                ]
            },
        },
    )
    try:
        events = [
            event
            async for event in runner.run(
                "publish the report",
                session_key,
                tool_context=tool_context,
                history_has_persisted_user=False,
                no_memory_capture=True,
            )
        ]
        return control_calls, qa_calls, events
    finally:
        await storage.close()


def _goal_artifact_topology_config(
    tmp_path,
    *,
    router: bool = False,
    ensemble: bool = False,
) -> GatewayConfig:
    return GatewayConfig(
        attachments=AttachmentsConfig(media_root=str(tmp_path / "media")),
        squilla_router=SquillaRouterConfig(enabled=router),
        llm_ensemble={
            "enabled": ensemble,
            "selection_mode": "router_dynamic",
        },
        agent_max_provider_retries=0,
    )


@pytest.mark.asyncio
async def test_goal_post_publish_router_decides_once_for_the_whole_normal_loop(
    tmp_path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    route_calls: list[str] = []

    async def route_once(turn: Any) -> Any:
        route_calls.append(turn.message)
        turn.model = "deepseek-v4-flash"
        turn.metadata.update(
            {
                "routed_tier": "c1",
                "routed_model": "deepseek-v4-flash",
                "baseline_model": "deepseek-v4-pro",
                "routing_source": "router",
                "routing_applied": True,
                "routing_confidence": 0.9,
            }
        )
        return turn

    route_once.__name__ = "apply_squilla_router"
    monkeypatch.setattr("openstarry_code.engine.steps.apply_squilla_router", route_once)
    provider = _GoalPostPublishLoopProvider()
    selector = _GoalArtifactTopologySelector(provider)

    control_calls, qa_calls, events = await _run_goal_artifact_topology(
        tmp_path,
        provider=provider,
        selector=selector,
        config=_goal_artifact_topology_config(tmp_path, router=True),
    )

    assert route_calls == ["publish the report"]
    assert provider.calls == 4
    assert len([event for event in events if isinstance(event, RouterDecisionEvent)]) == 1
    assert control_calls == ["progress:completed", "goal:complete"]
    assert qa_calls == []
    _assert_goal_artifact_published_once(events)
    assert all(
        "Goal delivery checkpoint"
        not in "\n".join(str(message.content) for message in request)
        for request in provider.requests
    )


@pytest.mark.asyncio
async def test_goal_post_publish_selector_keeps_the_active_fallback_leg(
    tmp_path,
) -> None:
    fallback = _GoalPostPublishLoopProvider()
    selector = _GoalArtifactFallbackSelector(fallback)

    control_calls, qa_calls, events = await _run_goal_artifact_topology(
        tmp_path,
        provider=fallback,
        selector=selector,
        config=_goal_artifact_topology_config(tmp_path),
    )

    assert selector.primary.calls == 1
    assert selector.fallback_calls == 1
    assert selector.current_config.model == "deepseek-v4-flash"
    assert fallback.calls == 4
    assert control_calls == ["progress:completed", "goal:complete"]
    assert qa_calls == []
    _assert_goal_artifact_published_once(events)
    assert all(
        set(tool_names)
        == {"publish_artifact", "qa_check", "update_goal", "update_goal_progress"}
        for tool_names in fallback.tool_names_seen[:3]
    )
    assert fallback.tool_names_seen[3] == []


@pytest.mark.asyncio
async def test_goal_post_publish_ensemble_runs_each_normal_decision_once(
    tmp_path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    provider = _GoalArtifactEnsembleProvider()
    selector = _GoalArtifactTopologySelector(provider)
    ensemble_builds: list[Any] = []

    def build_ensemble(**kwargs: Any) -> _GoalArtifactEnsembleProvider:
        ensemble_builds.append(kwargs["fallback_provider"])
        return provider

    monkeypatch.setattr(
        "openstarry_code.provider.ensemble.build_ensemble_provider_from_config",
        build_ensemble,
    )

    control_calls, qa_calls, events = await _run_goal_artifact_topology(
        tmp_path,
        provider=provider,
        selector=selector,
        config=_goal_artifact_topology_config(tmp_path, ensemble=True),
    )

    assert len(ensemble_builds) == 1
    assert provider.calls == 4
    assert control_calls == ["progress:completed", "goal:complete"]
    assert qa_calls == []
    _assert_goal_artifact_published_once(events)
    done = next(event for event in events if isinstance(event, DoneEvent))
    usage_by_role = {row["role"]: row for row in done.model_usage_breakdown}
    assert usage_by_role["proposer"]["request_count"] == 4
    assert usage_by_role["aggregator"]["request_count"] == 4
    assert done.ensemble_trace is not None
    assert done.ensemble_trace["llm_request_count"] == 8
