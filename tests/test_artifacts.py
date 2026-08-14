from __future__ import annotations

import hashlib
import io
import json
import os
import shutil
from dataclasses import replace
from pathlib import Path
from types import SimpleNamespace

import pytest
from pptx import Presentation

from openstarry_code.artifacts import (
    DEFAULT_ARTIFACT_DISK_BUDGET_BYTES,
    DEFAULT_ARTIFACT_MAX_BYTES,
    INSTALLER_ARTIFACT_MAX_BYTES,
    ArtifactBudgetError,
    ArtifactIntegrityError,
    ArtifactNotFoundError,
    ArtifactStore,
    artifact_cursor,
    artifact_marker,
    artifact_payload,
    strip_artifact_markers_from_text,
)
from openstarry_code.engine.types import ToolCall
from openstarry_code.session.plans import PlanRunConflictError
from openstarry_code.tools.builtin.artifacts import publish_artifact
from openstarry_code.tools.dispatch import build_tool_handler
from openstarry_code.tools.registry import ToolRegistry
from openstarry_code.tools.types import (
    CallerKind,
    RetryableToolInputError,
    ToolContext,
    ToolError,
    ToolSpec,
    current_tool_context,
)

PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


def _valid_pptx_bytes(title: str = "Validated deliverable") -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    assert slide.shapes.title is not None
    slide.shapes.title.text = title
    output = io.BytesIO()
    presentation.save(output)
    return output.getvalue()


def test_artifact_store_round_trips_metadata_and_bytes(tmp_path: Path) -> None:
    store = ArtifactStore(tmp_path)

    ref = store.publish_bytes(
        b"hello\n",
        session_id="session-1",
        session_key="agent:main:webchat:session-1",
        name="report.txt",
        mime="text/plain",
        source="publish_artifact",
    )
    path = store.path_for(ref)

    assert ref.kind == "artifact_ref"
    assert ref.name == "report.txt"
    assert ref.size == 6
    assert ref.download_url == "/api/v1/artifacts/" + ref.id
    assert path.read_bytes() == b"hello\n"

    resolved_ref, resolved_path = store.resolve_for_download(ref.id, session_id="session-1")
    assert resolved_ref == ref
    assert resolved_path == path


def _set_artifact_created_at(store: ArtifactStore, ref, created_at: str):
    updated = replace(ref, created_at=created_at)
    meta_path = store.path_for(ref).parent / "meta.json"
    meta_path.write_text(
        json.dumps(updated.to_dict(), ensure_ascii=False, sort_keys=True),
        encoding="utf-8",
    )
    return updated


def test_artifact_store_lists_stable_backwards_pages_and_gets_metadata(
    tmp_path: Path,
) -> None:
    store = ArtifactStore(tmp_path)
    session_id = "session-1"
    session_key = "agent:main:webchat:session-1"
    refs = [
        store.publish_bytes(
            f"payload-{index}".encode(),
            session_id=session_id,
            session_key=session_key,
            name=f"report-{index}.txt",
            mime="text/plain",
            source="publish_artifact",
        )
        for index in range(3)
    ]
    refs = [
        _set_artifact_created_at(store, ref, f"2026-01-0{index + 1}T00:00:00Z")
        for index, ref in enumerate(refs)
    ]

    newest = store.list_refs(session_id=session_id, limit=2)

    assert newest.refs == tuple(refs[1:])
    assert newest.has_more is True
    assert newest.total_count == 3
    assert store.get_ref(session_id=session_id, artifact_id=refs[2].id) == refs[2]

    older = store.list_refs(
        session_id=session_id,
        limit=2,
        before=artifact_cursor(newest.refs[0]),
    )

    assert older.refs == (refs[0],)
    assert older.has_more is False
    assert older.total_count == 3


def test_artifact_store_metadata_listing_skips_invalid_refs_without_reading_material(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    store = ArtifactStore(tmp_path)
    session_id = "session-1"
    session_key = "agent:main:webchat:session-1"
    good = store.publish_bytes(
        b"good",
        session_id=session_id,
        session_key=session_key,
        name="good.txt",
        mime="text/plain",
        source="publish_artifact",
    )
    corrupt = store.publish_bytes(
        b"corrupt metadata",
        session_id=session_id,
        session_key=session_key,
        name="corrupt.txt",
        mime="text/plain",
        source="publish_artifact",
    )
    non_object = store.publish_bytes(
        b"non-object metadata",
        session_id=session_id,
        session_key=session_key,
        name="non-object.txt",
        mime="text/plain",
        source="publish_artifact",
    )
    missing = store.publish_bytes(
        b"missing material",
        session_id=session_id,
        session_key=session_key,
        name="missing.txt",
        mime="text/plain",
        source="publish_artifact",
    )
    wrong_session = store.publish_bytes(
        b"wrong session metadata",
        session_id=session_id,
        session_key=session_key,
        name="wrong-session.txt",
        mime="text/plain",
        source="publish_artifact",
    )
    (store.path_for(corrupt).parent / "meta.json").write_text("{", encoding="utf-8")
    (store.path_for(non_object).parent / "meta.json").write_text("[]", encoding="utf-8")
    store.path_for(missing).unlink()
    wrong_session_meta = store.path_for(wrong_session).parent / "meta.json"
    wrong_session_meta.write_text(
        json.dumps(
            replace(wrong_session, session_id="session-2").to_dict(),
            ensure_ascii=False,
            sort_keys=True,
        ),
        encoding="utf-8",
    )

    def _unexpected_material_read(_path: Path) -> bytes:
        raise AssertionError("artifact metadata listing must not read material bytes")

    monkeypatch.setattr(Path, "read_bytes", _unexpected_material_read)

    page = store.list_refs(session_id=session_id, limit=20)

    assert page.refs == (good,)
    assert page.total_count == 1
    assert store.get_ref(session_id=session_id, artifact_id=good.id) == good
    for invalid_ref in (corrupt, non_object, missing):
        with pytest.raises(ArtifactNotFoundError):
            store.get_ref(session_id=session_id, artifact_id=invalid_ref.id)
    with pytest.raises(ArtifactNotFoundError):
        store.get_ref(session_id="session-2", artifact_id=good.id)


@pytest.mark.parametrize("invalid_field", ["id", "session_id"])
def test_artifact_store_validates_metadata_identity_before_material_lookup(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
    invalid_field: str,
) -> None:
    store = ArtifactStore(tmp_path)
    session_id = "session-1"
    ref = store.publish_bytes(
        b"material",
        session_id=session_id,
        session_key="agent:main:webchat:session-1",
        name="report.txt",
        mime="text/plain",
        source="publish_artifact",
    )
    invalid_ref = replace(
        ref,
        **(
            {"id": "art-other-session-artifact"}
            if invalid_field == "id"
            else {"session_id": "session-2"}
        ),
    )
    (store.path_for(ref).parent / "meta.json").write_text(
        json.dumps(invalid_ref.to_dict(), ensure_ascii=False, sort_keys=True),
        encoding="utf-8",
    )

    def _unexpected_material_lookup(_ref):
        raise PermissionError("another session must not be inspected")

    monkeypatch.setattr(
        store,
        "_preferred_material_path_for_ref",
        _unexpected_material_lookup,
    )

    page = store.list_refs(session_id=session_id, limit=10)

    assert page.refs == ()
    with pytest.raises(ArtifactNotFoundError):
        store.get_ref(session_id=session_id, artifact_id=ref.id)


def test_artifact_store_listing_propagates_directory_io_errors(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    store = ArtifactStore(tmp_path)

    def _unreadable_directory(_session_id: str):
        raise OSError("artifact directory is unreadable")

    monkeypatch.setattr(
        store,
        "_iter_session_meta_paths_for_listing",
        _unreadable_directory,
    )

    with pytest.raises(OSError, match="artifact directory is unreadable"):
        store.list_refs(session_id="session-1", limit=20)


def test_artifact_store_rejects_invalid_or_unknown_list_cursor(tmp_path: Path) -> None:
    store = ArtifactStore(tmp_path)

    with pytest.raises(ValueError, match="artifact id is invalid"):
        store.list_refs(session_id="session-1", limit=20, before="not-a-cursor")
    with pytest.raises(ValueError, match="artifact cursor not found"):
        store.list_refs(session_id="session-1", limit=20, before="art-missing")


@pytest.mark.parametrize("layout", ["current", "legacy-short", "legacy-plain"])
def test_artifact_store_list_reads_every_supported_layout(
    tmp_path: Path,
    layout: str,
) -> None:
    store = ArtifactStore(tmp_path)
    session_id = "532d5065-abce-499f-97b0-bbf2a067d5ab"
    ref = store.publish_bytes(
        b"layout material",
        session_id=session_id,
        session_key="agent:main:webchat:layout",
        name="layout.txt",
        mime="text/plain",
        source="publish_artifact",
    )
    current_dir = store.path_for(ref).parent
    if layout == "legacy-short":
        session_token = hashlib.sha256(session_id.encode("utf-8")).hexdigest()[:16]
        artifact_token = hashlib.sha256(ref.id.encode("utf-8")).hexdigest()[:16]
        target_dir = tmp_path / "artifacts" / "s" / session_token / artifact_token
        target_dir.parent.mkdir(parents=True)
        current_dir.rename(target_dir)
    elif layout == "legacy-plain":
        from openstarry_code.artifacts import _safe_token

        target_dir = tmp_path / "artifacts" / _safe_token(session_id) / ref.id
        target_dir.mkdir(parents=True)
        (current_dir / "data").rename(target_dir / ref.sha256)
        (current_dir / "meta.json").rename(target_dir / "meta.json")

    page = store.list_refs(session_id=session_id, limit=10)

    assert page.refs == (ref,)
    assert store.get_ref(session_id=session_id, artifact_id=ref.id) == ref


def test_artifact_store_list_prefers_current_duplicate_and_sorts_equal_times(
    tmp_path: Path,
) -> None:
    store = ArtifactStore(tmp_path)
    session_id = "532d5065-abce-499f-97b0-bbf2a067d5ab"
    session_key = "agent:main:webchat:duplicates"
    first = store.publish_bytes(
        b"first",
        session_id=session_id,
        session_key=session_key,
        name="current.txt",
        mime="text/plain",
        source="publish_artifact",
    )
    second = store.publish_bytes(
        b"second",
        session_id=session_id,
        session_key=session_key,
        name="second.txt",
        mime="text/plain",
        source="publish_artifact",
    )
    first_dir = store.path_for(first).parent
    legacy_session_token = hashlib.sha256(session_id.encode("utf-8")).hexdigest()[:16]
    legacy_artifact_token = hashlib.sha256(first.id.encode("utf-8")).hexdigest()[:16]
    legacy_dir = tmp_path / "artifacts" / "s" / legacy_session_token / legacy_artifact_token
    shutil.copytree(first_dir, legacy_dir)
    legacy_ref = replace(first, name="legacy-duplicate.txt")
    (legacy_dir / "meta.json").write_text(
        json.dumps(legacy_ref.to_dict(), ensure_ascii=False, sort_keys=True),
        encoding="utf-8",
    )
    same_time = "2026-03-01T00:00:00Z"
    first = _set_artifact_created_at(store, first, same_time)
    second = _set_artifact_created_at(store, second, same_time)

    page = store.list_refs(session_id=session_id, limit=10)

    assert [ref.id for ref in page.refs] == sorted([first.id, second.id])
    assert page.total_count == 2
    selected_first = next(ref for ref in page.refs if ref.id == first.id)
    assert selected_first.name == "current.txt"


def test_artifact_store_list_uses_preferred_meta_with_split_material_layout(
    tmp_path: Path,
) -> None:
    store = ArtifactStore(tmp_path)
    session_id = "532d5065-abce-499f-97b0-bbf2a067d5ab"
    ref = store.publish_bytes(
        b"current material",
        session_id=session_id,
        session_key="agent:main:webchat:mixed-layout",
        name="current.txt",
        mime="text/plain",
        source="publish_artifact",
    )
    current_dir = store.path_for(ref).parent
    legacy_session_token = hashlib.sha256(session_id.encode("utf-8")).hexdigest()[:16]
    legacy_artifact_token = hashlib.sha256(ref.id.encode("utf-8")).hexdigest()[:16]
    legacy_dir = tmp_path / "artifacts" / "s" / legacy_session_token / legacy_artifact_token
    shutil.copytree(current_dir, legacy_dir)

    legacy_bytes = b"legacy material"
    legacy_ref = replace(
        ref,
        sha256=hashlib.sha256(legacy_bytes).hexdigest(),
        name="legacy.txt",
        size=len(legacy_bytes),
    )
    (legacy_dir / "data").write_bytes(legacy_bytes)
    (legacy_dir / "meta.json").write_text(
        json.dumps(legacy_ref.to_dict(), ensure_ascii=False, sort_keys=True),
        encoding="utf-8",
    )
    (current_dir / "data").unlink()

    page = store.list_refs(session_id=session_id, limit=10)

    assert page.refs == (ref,)
    assert store.get_ref(session_id=session_id, artifact_id=ref.id) == ref
    with pytest.raises(ArtifactIntegrityError):
        store.resolve_for_download(ref.id, session_id=session_id)


def test_artifact_store_split_layout_with_matching_material_remains_downloadable(
    tmp_path: Path,
) -> None:
    store = ArtifactStore(tmp_path)
    session_id = "532d5065-abce-499f-97b0-bbf2a067d5ab"
    ref = store.publish_bytes(
        b"shared material",
        session_id=session_id,
        session_key="agent:main:webchat:split-layout",
        name="current.txt",
        mime="text/plain",
        source="publish_artifact",
    )
    current_dir = store.path_for(ref).parent
    legacy_session_token = hashlib.sha256(session_id.encode("utf-8")).hexdigest()[:16]
    legacy_artifact_token = hashlib.sha256(ref.id.encode("utf-8")).hexdigest()[:16]
    legacy_dir = tmp_path / "artifacts" / "s" / legacy_session_token / legacy_artifact_token
    shutil.copytree(current_dir, legacy_dir)
    (current_dir / "data").unlink()

    page = store.list_refs(session_id=session_id, limit=10)
    fetched = store.get_ref(session_id=session_id, artifact_id=ref.id)
    resolved_ref, resolved_path = store.resolve_for_download(ref.id, session_id=session_id)

    assert page.refs == (ref,)
    assert fetched == ref
    assert resolved_ref == ref
    assert resolved_path == legacy_dir / "data"
    assert resolved_path.read_bytes() == b"shared material"


@pytest.mark.skipif(os.name == "nt", reason="POSIX symlink regression")
@pytest.mark.parametrize("component", ["root", "meta", "material"])
def test_artifact_store_list_and_get_never_follow_symlink_targets(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
    component: str,
) -> None:
    store = ArtifactStore(tmp_path)
    session_id = "session-links"
    ref = store.publish_bytes(
        b"material",
        session_id=session_id,
        session_key="agent:main:webchat:links",
        name="report.txt",
        mime="text/plain",
        source="publish_artifact",
    )
    artifact_dir = store.path_for(ref).parent
    root = store._artifact_session_roots(session_id)[0]
    selected_path = {
        "root": root,
        "meta": artifact_dir / "meta.json",
        "material": artifact_dir / "data",
    }[component]
    target = tmp_path / f"outside-{component}"
    selected_path.rename(target)
    selected_path.symlink_to(target, target_is_directory=component == "root")

    original_stat = Path.stat

    def _reject_follow(
        path: Path,
        *,
        follow_symlinks: bool = True,
    ):
        if path == selected_path and follow_symlinks:
            raise AssertionError("artifact listing followed a symlink target")
        return original_stat(path, follow_symlinks=follow_symlinks)

    monkeypatch.setattr(Path, "stat", _reject_follow)

    page = store.list_refs(session_id=session_id, limit=10)

    assert page.refs == ()
    with pytest.raises(ArtifactNotFoundError):
        store.get_ref(session_id=session_id, artifact_id=ref.id)


def test_artifact_copy_keeps_using_the_existing_metadata_iterator(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    store = ArtifactStore(tmp_path)
    scanned: list[str] = []

    def _existing_iterator(session_id: str):
        scanned.append(session_id)
        return iter(())

    def _listing_iterator(_session_id: str):
        raise AssertionError("fork copy must not use the list-only iterator")

    monkeypatch.setattr(store, "_iter_session_meta_paths", _existing_iterator)
    monkeypatch.setattr(
        store,
        "_iter_session_meta_paths_for_listing",
        _listing_iterator,
    )

    assert (
        store.copy_session_artifacts(
            source_session_id="source",
            target_session_id="target",
            target_session_key="agent:main:webchat:target",
        )
        == 0
    )
    assert scanned == ["source"]


def test_artifact_store_finds_existing_session_deliverable_by_name_and_sha(
    tmp_path: Path,
) -> None:
    store = ArtifactStore(tmp_path)
    ref = store.publish_bytes(
        b"pptx bytes",
        session_id="session-1",
        session_key="agent:main:webchat:session-1",
        name="brief.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        source="create_pptx",
    )

    found = store.find_existing_ref(
        session_id="session-1",
        session_key="agent:main:webchat:session-1",
        sha256=ref.sha256,
        name="brief.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )

    assert found == ref
    assert (
        store.find_existing_ref(
            session_id="session-2",
            session_key="agent:main:webchat:session-2",
            sha256=ref.sha256,
            name="brief.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
        is None
    )


def test_artifact_store_skips_existing_deliverable_with_bad_material(
    tmp_path: Path,
) -> None:
    store = ArtifactStore(tmp_path)
    ref = store.publish_bytes(
        b"pptx bytes",
        session_id="session-1",
        session_key="agent:main:webchat:session-1",
        name="brief.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        source="create_pptx",
    )
    store.path_for(ref).write_bytes(b"corrupt")

    assert (
        store.find_existing_ref(
            session_id="session-1",
            session_key="agent:main:webchat:session-1",
            sha256=ref.sha256,
            name="brief.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
        is None
    )


def test_artifact_store_uses_short_material_paths_for_uuid_sessions(tmp_path: Path) -> None:
    store = ArtifactStore(tmp_path)
    long_root = tmp_path / ("deep-root-" + ("x" * 80))
    store = ArtifactStore(long_root)
    session_id = "532d5065-abce-499f-97b0-bbf2a067d5ab"

    ref = store.publish_bytes(
        b"pptx",
        session_id=session_id,
        session_key="agent:main:webchat:default",
        name="北京2027房价预测分析报告.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        source="publish_artifact",
    )

    material_path = store.path_for(ref)
    assert material_path.name == "data"
    assert session_id not in str(material_path)
    assert len(str(material_path)) < 260
    resolved_ref, resolved_path = store.resolve_for_download(ref.id, session_id=session_id)
    assert resolved_ref == ref
    assert resolved_path == material_path


def test_artifact_store_resolves_legacy_short_material_paths(tmp_path: Path) -> None:
    store = ArtifactStore(tmp_path)
    session_id = "532d5065-abce-499f-97b0-bbf2a067d5ab"

    ref = store.publish_bytes(
        b"pptx",
        session_id=session_id,
        session_key="agent:main:webchat:default",
        name="brief.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        source="publish_artifact",
    )

    current_dir = store.path_for(ref).parent
    legacy_session_token = hashlib.sha256(session_id.encode("utf-8")).hexdigest()[:16]
    legacy_artifact_token = hashlib.sha256(ref.id.encode("utf-8")).hexdigest()[:16]
    legacy_dir = tmp_path / "artifacts" / "s" / legacy_session_token / legacy_artifact_token
    legacy_dir.parent.mkdir(parents=True)
    current_dir.rename(legacy_dir)

    resolved_ref, resolved_path = store.resolve_for_download(ref.id, session_id=session_id)

    assert resolved_ref == ref
    assert resolved_path == legacy_dir / "data"


def test_artifact_store_resolves_legacy_short_thumbnail_paths(tmp_path: Path) -> None:
    from PIL import Image

    store = ArtifactStore(tmp_path)
    session_id = "532d5065-abce-499f-97b0-bbf2a067d5ab"
    out = io.BytesIO()
    Image.new("RGB", (8, 8), color="red").save(out, format="PNG")

    ref = store.publish_bytes(
        out.getvalue(),
        session_id=session_id,
        session_key="agent:main:webchat:default",
        name="chart.png",
        mime="image/png",
        source="publish_artifact",
    )
    assert ref.has_thumbnail is True

    current_dir = store.path_for(ref).parent
    legacy_session_token = hashlib.sha256(session_id.encode("utf-8")).hexdigest()[:16]
    legacy_artifact_token = hashlib.sha256(ref.id.encode("utf-8")).hexdigest()[:16]
    legacy_dir = tmp_path / "artifacts" / "s" / legacy_session_token / legacy_artifact_token
    legacy_dir.parent.mkdir(parents=True)
    current_dir.rename(legacy_dir)

    thumbnail = store.resolve_thumbnail_for_download(ref.id, session_id=session_id)

    assert thumbnail is not None
    resolved_ref, thumbnail_path = thumbnail
    assert resolved_ref == ref
    assert thumbnail_path == legacy_dir / "thumb.webp"


def test_artifact_payload_omits_session_key_and_query_token(tmp_path: Path) -> None:
    store = ArtifactStore(tmp_path)
    ref = store.publish_bytes(
        b"hello\n",
        session_id="session-1",
        session_key="agent:main:webchat:session-1",
        name="report.txt",
        mime="text/plain",
        source="publish_artifact",
    )

    payload = artifact_payload(ref)

    assert "session_key" not in payload
    assert "sessionKey" not in json.dumps(payload)
    assert payload["download_url"] == f"/api/v1/artifacts/{ref.id}"


def test_artifact_payload_keeps_thumbnail_url_across_persist_and_replay() -> None:
    # Live event carries the internal has_thumbnail boolean; the public payload
    # exposes only the reconstructed thumbnail_url string.
    live = artifact_payload(
        SimpleNamespace(
            id="art-bmYMIceM2Ddx3rkFM4BOmZ7A",
            kind="artifact_ref",
            sha256="a" * 64,
            name="chart.png",
            mime="image/png",
            size=954199,
            session_id="session-1",
            source="publish_artifact",
            created_at="2026-06-13T00:00:00Z",
            store="artifacts",
            download_url="/api/v1/artifacts/art-bmYMIceM2Ddx3rkFM4BOmZ7A",
            has_thumbnail=True,
        )
    )
    assert "has_thumbnail" not in live
    assert live["thumbnail_url"] == "/api/v1/artifacts/art-bmYMIceM2Ddx3rkFM4BOmZ7A?variant=thumb"

    # Replaying the persisted public payload (which no longer carries the boolean)
    # must rebuild the same thumbnail_url instead of falling back to the full file.
    persisted = json.loads(json.dumps(live))
    replayed = artifact_payload(persisted)
    assert replayed["thumbnail_url"] == live["thumbnail_url"]


def test_artifact_payload_omits_thumbnail_url_without_thumbnail() -> None:
    no_thumb = artifact_payload(
        SimpleNamespace(
            id="art-NoThumbXXXXXXXXXXXXXXXXX",
            kind="artifact_ref",
            sha256="b" * 64,
            name="doc.pdf",
            mime="application/pdf",
            size=1000,
            session_id="session-1",
            source="publish_artifact",
            created_at="2026-06-13T00:00:00Z",
            store="artifacts",
            download_url="/api/v1/artifacts/art-NoThumbXXXXXXXXXXXXXXXXX",
            has_thumbnail=False,
        )
    )
    assert "thumbnail_url" not in no_thumb
    replayed = artifact_payload(json.loads(json.dumps(no_thumb)))
    assert "thumbnail_url" not in replayed


def test_artifact_store_preserves_unicode_filename_and_normalizes_mime_params(
    tmp_path: Path,
) -> None:
    store = ArtifactStore(tmp_path)

    ref = store.publish_bytes(
        b"hello\n",
        session_id="session-1",
        session_key="agent:main:webchat:session-1",
        name="记忆修补师.txt",
        mime="text/plain; charset=utf-8",
        source="publish_artifact",
    )

    assert ref.name == "记忆修补师.txt"
    assert ref.mime == "text/plain"


def test_artifact_store_rejects_hash_mismatch(tmp_path: Path) -> None:
    store = ArtifactStore(tmp_path)
    ref = store.publish_bytes(
        b"hello",
        session_id="session-1",
        session_key="agent:main:webchat:session-1",
        name="report.txt",
        mime="text/plain",
        source="publish_artifact",
    )

    store.path_for(ref).write_bytes(b"tampered")

    with pytest.raises(ArtifactIntegrityError):
        store.resolve_for_download(ref.id, session_id="session-1")


def test_artifact_store_enforces_per_file_and_disk_budgets(tmp_path: Path) -> None:
    store = ArtifactStore(tmp_path)

    with pytest.raises(ArtifactBudgetError):
        store.publish_bytes(
            b"abcdef",
            session_id="session-1",
            session_key="agent:main:webchat:session-1",
            name="too-big.txt",
            mime="text/plain",
            source="publish_artifact",
            max_bytes=5,
        )

    assert not list((tmp_path / "artifacts").rglob("too-big.txt"))

    store.publish_bytes(
        b"abc",
        session_id="session-1",
        session_key="agent:main:webchat:session-1",
        name="ok.txt",
        mime="text/plain",
        source="publish_artifact",
        disk_budget_bytes=6,
    )
    with pytest.raises(ArtifactBudgetError):
        store.publish_bytes(
            b"defg",
            session_id="session-1",
            session_key="agent:main:webchat:session-1",
            name="over-budget.txt",
            mime="text/plain",
            source="publish_artifact",
            disk_budget_bytes=6,
        )


def test_artifact_budget_defaults_are_open_source_sized() -> None:
    assert DEFAULT_ARTIFACT_MAX_BYTES == 30 * 1024 * 1024
    assert DEFAULT_ARTIFACT_DISK_BUDGET_BYTES == 512 * 1024 * 1024
    assert INSTALLER_ARTIFACT_MAX_BYTES == DEFAULT_ARTIFACT_DISK_BUDGET_BYTES


@pytest.mark.asyncio
async def test_publish_artifact_tool_allows_workspace_file_only(tmp_path: Path) -> None:
    workspace = tmp_path / "workspace"
    workspace.mkdir()
    output = workspace / "report.txt"
    output.write_text("ready", encoding="utf-8")
    ctx = ToolContext(
        is_owner=True,
        caller_kind=CallerKind.WEB,
        workspace_dir=str(workspace),
        artifact_media_root=str(tmp_path / "media"),
        artifact_session_id="session-1",
        session_key="agent:main:webchat:session-1",
    )

    token = current_tool_context.set(ctx)
    try:
        result = await publish_artifact(path="report.txt", name="final.txt", mime="text/plain")
    finally:
        current_tool_context.reset(token)

    payload = json.loads(result)
    assert payload["status"] == "published"
    assert payload["artifact"]["name"] == "final.txt"
    assert payload["artifact"]["mime"] == "text/plain"
    assert payload["artifact"]["session_id"] == "session-1"
    assert "session_key" not in payload["artifact"]
    assert "sessionKey" not in json.dumps(payload["artifact"])
    # The LLM-facing artifact has no URL — models tend to fabricate a host
    # when shown a relative URL ending in /api/v1/artifacts/...
    assert "download_url" not in payload["artifact"]
    assert payload["artifact"]["workspace_path"] == "report.txt"
    assert payload["artifact"]["local_path"] == str(output.resolve())
    assert "note" in payload
    assert "local_path" in payload["note"]
    assert "final response" in payload["note"]
    assert "Do not run more tools" in payload["note"]
    # The frontend event path still gets the full payload (with download_url).
    assert len(ctx.published_artifacts) == 1
    full_artifact = ctx.published_artifacts[0]
    assert full_artifact["download_url"] == f"/api/v1/artifacts/{full_artifact['id']}"
    llm_artifact = {
        k: v
        for k, v in payload["artifact"].items()
        if k not in {"workspace_path", "local_path"}
    }
    assert {k: v for k, v in full_artifact.items() if k != "download_url"} == llm_artifact


@pytest.mark.asyncio
async def test_publish_artifact_requires_completed_attached_plan_run(
    tmp_path: Path,
) -> None:
    class PlanStorage:
        async def get_plan_run(self, run_id: str) -> SimpleNamespace:
            assert run_id == "run-1"
            return SimpleNamespace(
                status="running",
                current_step_id="verify",
                active_task_id="task-1",
            )

    ctx = ToolContext(
        workspace_dir=str(tmp_path),
        artifact_media_root=str(tmp_path / "media"),
        artifact_session_id="session-1",
        session_key="agent:main:webchat:session-1",
        task_id="task-1",
        plan_run_id="run-1",
        plan_storage=PlanStorage(),
    )

    token = current_tool_context.set(ctx)
    try:
        with pytest.raises(RetryableToolInputError) as exc_info:
            await publish_artifact(path="report.txt")
    finally:
        current_tool_context.reset(token)

    message = exc_info.value.user_message
    assert "was not executed" in message
    assert "current step is verify" in message
    assert "after the final checkpoint returns no current step" in message
    assert ctx.published_artifacts == []


@pytest.mark.asyncio
async def test_publish_artifact_checkpoints_the_only_unfinished_final_step(
    tmp_path: Path,
) -> None:
    workspace = tmp_path / "workspace"
    workspace.mkdir()
    (workspace / "report.txt").write_text("ready", encoding="utf-8")

    class PlanStorage:
        def __init__(self) -> None:
            self.run = SimpleNamespace(
                status="running",
                current_step_id="publish",
                active_task_id="task-1",
                state_revision=7,
                step_states=[
                    {"step_id": "build", "status": "completed"},
                    {"step_id": "publish", "status": "in_progress"},
                ],
            )
            self.checkpoints: list[dict[str, object]] = []

        async def get_plan_run(self, run_id: str) -> SimpleNamespace:
            assert run_id == "run-1"
            return self.run

        async def checkpoint_plan_run(self, run_id: str, **kwargs: object) -> SimpleNamespace:
            assert run_id == "run-1"
            self.checkpoints.append(kwargs)
            self.run = SimpleNamespace(
                status="running",
                current_step_id=None,
                active_task_id="task-1",
                state_revision=8,
                step_states=[
                    {"step_id": "build", "status": "completed"},
                    {"step_id": "publish", "status": "completed"},
                ],
            )
            return self.run

    storage = PlanStorage()
    ctx = ToolContext(
        workspace_dir=str(workspace),
        artifact_media_root=str(tmp_path / "media"),
        artifact_session_id="session-1",
        session_key="agent:main:webchat:session-1",
        task_id="task-1",
        plan_run_id="run-1",
        plan_storage=storage,
    )

    token = current_tool_context.set(ctx)
    try:
        payload = json.loads(await publish_artifact(path="report.txt"))
    finally:
        current_tool_context.reset(token)

    assert payload["status"] == "published"
    assert storage.checkpoints == [
        {
            "expected_state_revision": 7,
            "step_id": "publish",
            "step_status": "completed",
            "next_step_id": None,
            "expected_active_task_id": "task-1",
        }
    ]
    assert len(ctx.published_artifacts) == 1


@pytest.mark.asyncio
@pytest.mark.parametrize(
    ("refreshed_status", "refreshed_active_task_id"),
    [
        ("cancelled", None),
        ("running", "task-2"),
    ],
)
async def test_publish_artifact_rejects_disallowed_state_after_checkpoint_conflict(
    tmp_path: Path,
    refreshed_status: str,
    refreshed_active_task_id: str | None,
) -> None:
    workspace = tmp_path / "workspace"
    workspace.mkdir()
    (workspace / "report.txt").write_text("ready", encoding="utf-8")

    class PlanStorage:
        def __init__(self) -> None:
            self.refreshed = False

        async def get_plan_run(self, run_id: str) -> SimpleNamespace:
            assert run_id == "run-1"
            if not self.refreshed:
                return SimpleNamespace(
                    status="running",
                    current_step_id="publish",
                    active_task_id="task-1",
                    state_revision=7,
                    step_states=[
                        {"step_id": "build", "status": "completed"},
                        {"step_id": "publish", "status": "in_progress"},
                    ],
                )
            return SimpleNamespace(
                status=refreshed_status,
                current_step_id=None,
                active_task_id=refreshed_active_task_id,
                state_revision=9,
                step_states=[
                    {"step_id": "build", "status": "completed"},
                    {"step_id": "publish", "status": "completed"},
                ],
            )

        async def checkpoint_plan_run(self, run_id: str, **kwargs: object) -> None:
            assert run_id == "run-1"
            self.refreshed = True
            raise PlanRunConflictError("plan run changed")

    ctx = ToolContext(
        workspace_dir=str(workspace),
        artifact_media_root=str(tmp_path / "media"),
        artifact_session_id="session-1",
        session_key="agent:main:webchat:session-1",
        task_id="task-1",
        plan_run_id="run-1",
        plan_storage=PlanStorage(),
    )

    token = current_tool_context.set(ctx)
    try:
        with pytest.raises(RetryableToolInputError, match="changed"):
            await publish_artifact(path="report.txt")
    finally:
        current_tool_context.reset(token)

    assert ctx.published_artifacts == []


@pytest.mark.asyncio
async def test_publish_artifact_allows_owned_delivery_after_checkpoint_conflict(
    tmp_path: Path,
) -> None:
    workspace = tmp_path / "workspace"
    workspace.mkdir()
    (workspace / "report.txt").write_text("ready", encoding="utf-8")

    class PlanStorage:
        def __init__(self) -> None:
            self.refreshed = False

        async def get_plan_run(self, run_id: str) -> SimpleNamespace:
            assert run_id == "run-1"
            if not self.refreshed:
                return SimpleNamespace(
                    status="running",
                    current_step_id="publish",
                    active_task_id="task-1",
                    state_revision=7,
                    step_states=[
                        {"step_id": "build", "status": "completed"},
                        {"step_id": "publish", "status": "in_progress"},
                    ],
                )
            return SimpleNamespace(
                status="running",
                current_step_id=None,
                active_task_id="task-1",
                state_revision=8,
                step_states=[
                    {"step_id": "build", "status": "completed"},
                    {"step_id": "publish", "status": "completed"},
                ],
            )

        async def checkpoint_plan_run(self, run_id: str, **kwargs: object) -> None:
            assert run_id == "run-1"
            self.refreshed = True
            raise PlanRunConflictError("plan run changed")

    ctx = ToolContext(
        workspace_dir=str(workspace),
        artifact_media_root=str(tmp_path / "media"),
        artifact_session_id="session-1",
        session_key="agent:main:webchat:session-1",
        task_id="task-1",
        plan_run_id="run-1",
        plan_storage=PlanStorage(),
    )

    token = current_tool_context.set(ctx)
    try:
        payload = json.loads(await publish_artifact(path="report.txt"))
    finally:
        current_tool_context.reset(token)

    assert payload["status"] == "published"
    assert len(ctx.published_artifacts) == 1


@pytest.mark.asyncio
async def test_publish_artifact_does_not_checkpoint_before_the_final_step(
    tmp_path: Path,
) -> None:
    class PlanStorage:
        def __init__(self) -> None:
            self.checkpointed = False

        async def get_plan_run(self, run_id: str) -> SimpleNamespace:
            assert run_id == "run-1"
            return SimpleNamespace(
                status="running",
                current_step_id="build",
                active_task_id="task-1",
                state_revision=3,
                step_states=[
                    {"step_id": "build", "status": "in_progress"},
                    {"step_id": "publish", "status": "pending"},
                ],
            )

        async def checkpoint_plan_run(self, run_id: str, **kwargs: object) -> None:
            self.checkpointed = True

    storage = PlanStorage()
    ctx = ToolContext(
        workspace_dir=str(tmp_path),
        artifact_media_root=str(tmp_path / "media"),
        artifact_session_id="session-1",
        session_key="agent:main:webchat:session-1",
        task_id="task-1",
        plan_run_id="run-1",
        plan_storage=storage,
    )

    token = current_tool_context.set(ctx)
    try:
        with pytest.raises(RetryableToolInputError):
            await publish_artifact(path="report.txt")
    finally:
        current_tool_context.reset(token)

    assert storage.checkpointed is False
    assert ctx.published_artifacts == []


@pytest.mark.asyncio
async def test_publish_artifact_does_not_checkpoint_an_invalid_final_artifact(
    tmp_path: Path,
) -> None:
    class PlanStorage:
        def __init__(self) -> None:
            self.checkpointed = False

        async def get_plan_run(self, run_id: str) -> SimpleNamespace:
            assert run_id == "run-1"
            return SimpleNamespace(
                status="running",
                current_step_id="publish",
                active_task_id="task-1",
                state_revision=4,
                step_states=[
                    {"step_id": "build", "status": "completed"},
                    {"step_id": "publish", "status": "in_progress"},
                ],
            )

        async def checkpoint_plan_run(self, run_id: str, **kwargs: object) -> None:
            self.checkpointed = True

    storage = PlanStorage()
    ctx = ToolContext(
        workspace_dir=str(tmp_path),
        artifact_media_root=str(tmp_path / "media"),
        artifact_session_id="session-1",
        session_key="agent:main:webchat:session-1",
        task_id="task-1",
        plan_run_id="run-1",
        plan_storage=storage,
    )

    token = current_tool_context.set(ctx)
    try:
        with pytest.raises(ToolError, match="artifact file not found"):
            await publish_artifact(path="missing.txt")
    finally:
        current_tool_context.reset(token)

    assert storage.checkpointed is False
    assert ctx.published_artifacts == []


@pytest.mark.asyncio
async def test_publish_artifact_allows_completed_attached_plan_run(
    tmp_path: Path,
) -> None:
    workspace = tmp_path / "workspace"
    workspace.mkdir()
    (workspace / "report.txt").write_text("ready", encoding="utf-8")

    class PlanStorage:
        async def get_plan_run(self, run_id: str) -> SimpleNamespace:
            assert run_id == "run-1"
            return SimpleNamespace(status="completed", current_step_id=None)

    ctx = ToolContext(
        workspace_dir=str(workspace),
        artifact_media_root=str(tmp_path / "media"),
        artifact_session_id="session-1",
        session_key="agent:main:webchat:session-1",
        plan_run_id="run-1",
        plan_storage=PlanStorage(),
    )

    token = current_tool_context.set(ctx)
    try:
        payload = json.loads(await publish_artifact(path="report.txt"))
    finally:
        current_tool_context.reset(token)

    assert payload["status"] == "published"
    assert len(ctx.published_artifacts) == 1


@pytest.mark.asyncio
async def test_publish_artifact_allows_delivery_ready_running_plan_run(
    tmp_path: Path,
) -> None:
    workspace = tmp_path / "workspace"
    workspace.mkdir()
    (workspace / "report.txt").write_text("ready", encoding="utf-8")

    class PlanStorage:
        async def get_plan_run(self, run_id: str) -> SimpleNamespace:
            assert run_id == "run-1"
            return SimpleNamespace(
                status="running",
                current_step_id=None,
                active_task_id="task-1",
                step_states=[
                    {"step_id": "build", "status": "completed"},
                    {"step_id": "verify", "status": "skipped"},
                ],
            )

    ctx = ToolContext(
        workspace_dir=str(workspace),
        artifact_media_root=str(tmp_path / "media"),
        artifact_session_id="session-1",
        session_key="agent:main:webchat:session-1",
        task_id="task-1",
        plan_run_id="run-1",
        plan_storage=PlanStorage(),
    )

    token = current_tool_context.set(ctx)
    try:
        payload = json.loads(await publish_artifact(path="report.txt"))
    finally:
        current_tool_context.reset(token)

    assert payload["status"] == "published"
    assert len(ctx.published_artifacts) == 1


@pytest.mark.asyncio
async def test_publish_artifact_rejects_paused_plan_run_without_retry(
    tmp_path: Path,
) -> None:
    class PlanStorage:
        async def get_plan_run(self, run_id: str) -> SimpleNamespace:
            assert run_id == "run-1"
            return SimpleNamespace(
                status="paused",
                current_step_id=None,
                step_states=[{"step_id": "build", "status": "completed"}],
            )

    ctx = ToolContext(
        workspace_dir=str(tmp_path),
        artifact_media_root=str(tmp_path / "media"),
        artifact_session_id="session-1",
        session_key="agent:main:webchat:session-1",
        plan_run_id="run-1",
        plan_storage=PlanStorage(),
    )

    token = current_tool_context.set(ctx)
    try:
        with pytest.raises(ToolError) as exc_info:
            await publish_artifact(path="report.txt")
    finally:
        current_tool_context.reset(token)

    assert not isinstance(exc_info.value, RetryableToolInputError)
    assert "unavailable for this terminal or unowned PlanRun state" in str(
        exc_info.value
    )


@pytest.mark.asyncio
async def test_publish_artifact_tool_allows_large_installer_artifact(tmp_path: Path) -> None:
    workspace = tmp_path / "workspace"
    workspace.mkdir()
    output = workspace / "OpenStarry Code-0.4.0-arm64.dmg"
    with output.open("wb") as handle:
        handle.seek(DEFAULT_ARTIFACT_MAX_BYTES + 1)
        handle.write(b"x")
    ctx = ToolContext(
        is_owner=True,
        caller_kind=CallerKind.WEB,
        workspace_dir=str(workspace),
        artifact_media_root=str(tmp_path / "media"),
        artifact_session_id="session-1",
        session_key="agent:main:webchat:session-1",
    )

    token = current_tool_context.set(ctx)
    try:
        result = await publish_artifact(path=output.name)
    finally:
        current_tool_context.reset(token)

    payload = json.loads(result)
    assert payload["status"] == "published"
    assert payload["artifact"]["name"] == output.name
    assert payload["artifact"]["size"] > DEFAULT_ARTIFACT_MAX_BYTES
    assert payload["artifact"]["mime"] == "application/x-apple-diskimage"


@pytest.mark.asyncio
async def test_publish_artifact_tool_preserves_source_extension_for_display_name(
    tmp_path: Path,
) -> None:
    workspace = tmp_path / "workspace"
    workspace.mkdir()
    output = workspace / "generated-chart.png"
    output.write_bytes(b"\x89PNG\r\n\x1a\nimage bytes")
    ctx = ToolContext(
        is_owner=True,
        caller_kind=CallerKind.WEB,
        workspace_dir=str(workspace),
        artifact_media_root=str(tmp_path / "media"),
        artifact_session_id="session-1",
        session_key="agent:main:webchat:session-1",
    )

    token = current_tool_context.set(ctx)
    try:
        result = await publish_artifact(
            path="generated-chart.png",
            name="Friendly Chart",
        )
    finally:
        current_tool_context.reset(token)

    payload = json.loads(result)

    assert payload["status"] == "published"
    assert payload["artifact"]["name"] == "Friendly Chart.png"
    assert payload["artifact"]["mime"] == "image/png"
    assert ctx.published_artifacts[0]["name"] == "Friendly Chart.png"
    assert ctx.published_artifacts[0]["mime"] == "image/png"


@pytest.mark.asyncio
async def test_publish_artifact_tool_keeps_download_name_mime_when_source_is_generic(
    tmp_path: Path,
) -> None:
    workspace = tmp_path / "workspace"
    workspace.mkdir()
    output = workspace / "payload.bin"
    output.write_bytes(b"image bytes")
    ctx = ToolContext(
        is_owner=True,
        caller_kind=CallerKind.WEB,
        workspace_dir=str(workspace),
        artifact_media_root=str(tmp_path / "media"),
        artifact_session_id="session-1",
        session_key="agent:main:webchat:session-1",
    )

    token = current_tool_context.set(ctx)
    try:
        result = await publish_artifact(
            path="payload.bin",
            name="Friendly Chart.png",
        )
    finally:
        current_tool_context.reset(token)

    payload = json.loads(result)

    assert payload["artifact"]["name"] == "Friendly Chart.png"
    assert payload["artifact"]["mime"] == "image/png"


@pytest.mark.parametrize(
    ("source_name", "artifact_name", "mime"),
    [
        ("broken.PPTX", "download.bin", "application/octet-stream"),
        ("broken.bin", "download.PPTX", "application/octet-stream"),
        ("broken.bin", "download.bin", PPTX_MIME),
    ],
)
@pytest.mark.asyncio
async def test_publish_artifact_tool_rejects_invalid_pptx_from_any_format_signal(
    tmp_path: Path,
    source_name: str,
    artifact_name: str,
    mime: str,
) -> None:
    workspace = tmp_path / "workspace"
    workspace.mkdir()
    output = workspace / source_name
    output.write_bytes(b"not a zip package")
    media_root = tmp_path / "media"
    ctx = ToolContext(
        workspace_dir=str(workspace),
        artifact_media_root=str(media_root),
        artifact_session_id="session-1",
        session_key="agent:main:webchat:session-1",
    )

    token = current_tool_context.set(ctx)
    try:
        with pytest.raises(RetryableToolInputError) as exc_info:
            await publish_artifact(path=source_name, name=artifact_name, mime=mime)
    finally:
        current_tool_context.reset(token)

    assert exc_info.value.user_message
    assert str(output.resolve()) not in exc_info.value.user_message
    assert not ctx.published_artifacts
    assert not list(media_root.rglob("meta.json"))


@pytest.mark.asyncio
async def test_publish_artifact_invalid_pptx_dispatch_is_safe_and_retryable(
    tmp_path: Path,
) -> None:
    workspace = tmp_path / "workspace"
    workspace.mkdir()
    output = workspace / "brief.pptx"
    output.write_bytes(b"not a zip package")
    media_root = tmp_path / "media"
    ctx = ToolContext(
        workspace_dir=str(workspace),
        artifact_media_root=str(media_root),
        artifact_session_id="session-1",
        session_key="agent:main:webchat:session-1",
    )
    registry = ToolRegistry()
    registry.register(
        ToolSpec(
            name="publish_artifact",
            description="Publish a generated artifact.",
            parameters={
                "path": {"type": "string"},
                "name": {"type": "string"},
                "mime": {"type": "string"},
            },
            required=["path"],
        ),
        publish_artifact,
    )
    handler = build_tool_handler(registry, ctx)

    result = await handler(
        ToolCall(
            tool_use_id="tc-invalid-pptx",
            tool_name="publish_artifact",
            arguments={"path": output.name, "mime": PPTX_MIME},
        )
    )

    envelope = json.loads(result.content)
    assert result.is_error is True
    assert set(envelope) == {
        "status",
        "tool",
        "error_class",
        "user_message",
        "retry_allowed",
    }
    assert envelope["status"] == "error"
    assert envelope["tool"] == "publish_artifact"
    assert envelope["error_class"] == "RetryableToolInputError"
    assert envelope["retry_allowed"] is True
    assert envelope["user_message"]
    assert "regenerate" in envelope["user_message"].lower()
    assert "not a zip package" not in envelope["user_message"]
    assert str(output.resolve()) not in envelope["user_message"]
    assert result.artifacts == []
    assert not ctx.published_artifacts
    assert not list(media_root.rglob("meta.json"))


@pytest.mark.asyncio
async def test_publish_artifact_preflights_pptx_budget_before_read_or_validation(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    import openstarry_code.tools.builtin.artifacts as artifacts_module

    workspace = tmp_path / "workspace"
    workspace.mkdir()
    output = workspace / "brief.pptx"
    output.write_bytes(_valid_pptx_bytes("Validated before budget"))
    media_root = tmp_path / "media"
    ctx = ToolContext(
        workspace_dir=str(workspace),
        artifact_media_root=str(media_root),
        artifact_session_id="session-1",
        session_key="agent:main:webchat:session-1",
        artifact_max_bytes=4,
    )

    validation_calls = 0
    def unexpected_validation(*args: object, **kwargs: object) -> None:
        nonlocal validation_calls
        validation_calls += 1

    def unexpected_read(*args: object, **kwargs: object) -> None:
        pytest.fail("oversized PPTX must be rejected before read_bytes")

    monkeypatch.setattr(
        artifacts_module,
        "validate_artifact_for_delivery",
        unexpected_validation,
    )
    monkeypatch.setattr(artifacts_module.Path, "read_bytes", unexpected_read)
    token = current_tool_context.set(ctx)
    try:
        with pytest.raises(ToolError, match="artifact exceeds per-file budget"):
            await publish_artifact(path=output.name, mime=PPTX_MIME)
    finally:
        current_tool_context.reset(token)

    assert validation_calls == 0
    assert not ctx.published_artifacts
    assert not list(media_root.rglob("meta.json"))


@pytest.mark.asyncio
async def test_publish_artifact_rejects_oversized_pptx_before_historical_dedupe(
    tmp_path: Path,
) -> None:
    workspace = tmp_path / "workspace"
    workspace.mkdir()
    payload = _valid_pptx_bytes("Historical valid deck")
    output = workspace / "brief.pptx"
    output.write_bytes(payload)
    media_root = tmp_path / "media"
    store = ArtifactStore(media_root)
    store.publish_bytes(
        payload,
        session_id="session-1",
        session_key="agent:main:webchat:session-1",
        name=output.name,
        mime=PPTX_MIME,
        source="legacy",
        max_bytes=None,
    )
    ctx = ToolContext(
        workspace_dir=str(workspace),
        artifact_media_root=str(media_root),
        artifact_session_id="session-1",
        session_key="agent:main:webchat:session-1",
        artifact_max_bytes=4,
    )

    token = current_tool_context.set(ctx)
    try:
        with pytest.raises(ToolError, match="artifact exceeds per-file budget"):
            await publish_artifact(path=output.name, mime=PPTX_MIME)
    finally:
        current_tool_context.reset(token)

    assert ctx.published_artifacts == []


@pytest.mark.asyncio
async def test_publish_artifact_validates_pptx_before_current_turn_dedupe(
    tmp_path: Path,
) -> None:
    workspace = tmp_path / "workspace"
    workspace.mkdir()
    invalid_payload = b"not a zip package"
    output = workspace / "brief.pptx"
    output.write_bytes(invalid_payload)
    previous = {
        "id": "art-existing",
        "sha256": hashlib.sha256(invalid_payload).hexdigest(),
        "name": "brief.pptx",
        "mime": PPTX_MIME,
    }
    media_root = tmp_path / "media"
    ctx = ToolContext(
        workspace_dir=str(workspace),
        artifact_media_root=str(media_root),
        artifact_session_id="session-1",
        session_key="agent:main:webchat:session-1",
        published_artifacts=[previous],
    )

    token = current_tool_context.set(ctx)
    try:
        with pytest.raises(RetryableToolInputError):
            await publish_artifact(path=output.name, name=output.name, mime=PPTX_MIME)
    finally:
        current_tool_context.reset(token)

    assert ctx.published_artifacts == [previous]
    assert not list(media_root.rglob("meta.json"))


@pytest.mark.asyncio
async def test_publish_artifact_validates_pptx_before_persisted_dedupe(
    tmp_path: Path,
) -> None:
    workspace = tmp_path / "workspace"
    workspace.mkdir()
    invalid_payload = b"not a zip package"
    output = workspace / "brief.pptx"
    output.write_bytes(invalid_payload)
    media_root = tmp_path / "media"
    store = ArtifactStore(media_root)
    historical = store.publish_bytes(
        invalid_payload,
        session_id="session-1",
        session_key="agent:main:webchat:session-1",
        name=output.name,
        mime=PPTX_MIME,
        source="legacy_publish_artifact",
    )
    ctx = ToolContext(
        workspace_dir=str(workspace),
        artifact_media_root=str(media_root),
        artifact_session_id="session-1",
        session_key="agent:main:webchat:session-1",
    )

    token = current_tool_context.set(ctx)
    try:
        with pytest.raises(RetryableToolInputError):
            await publish_artifact(path=output.name, name=output.name, mime=PPTX_MIME)
    finally:
        current_tool_context.reset(token)

    assert not ctx.published_artifacts
    assert [path.parent for path in media_root.rglob("meta.json")] == [
        store.path_for(historical).parent
    ]


@pytest.mark.asyncio
async def test_publish_artifact_allows_valid_retry_after_invalid_pptx_is_replaced(
    tmp_path: Path,
) -> None:
    workspace = tmp_path / "workspace"
    workspace.mkdir()
    output = workspace / "brief.pptx"
    output.write_bytes(b"not a zip package")
    media_root = tmp_path / "media"
    ctx = ToolContext(
        workspace_dir=str(workspace),
        artifact_media_root=str(media_root),
        artifact_session_id="session-1",
        session_key="agent:main:webchat:session-1",
    )

    token = current_tool_context.set(ctx)
    try:
        with pytest.raises(RetryableToolInputError):
            await publish_artifact(path=output.name, mime=PPTX_MIME)
        valid_payload = _valid_pptx_bytes("Retry succeeded")
        output.write_bytes(valid_payload)
        result = json.loads(await publish_artifact(path=output.name, mime=PPTX_MIME))
    finally:
        current_tool_context.reset(token)

    assert result["status"] == "published"
    assert len(ctx.published_artifacts) == 1
    _, material_path = ArtifactStore(media_root).resolve_for_download(
        result["artifact"]["id"],
        session_id="session-1",
    )
    assert material_path.read_bytes() == valid_payload


@pytest.mark.asyncio
async def test_publish_artifact_stores_the_exact_pptx_bytes_that_were_validated(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    import openstarry_code.tools.builtin.artifacts as artifacts_module

    workspace = tmp_path / "workspace"
    workspace.mkdir()
    validated_payload = _valid_pptx_bytes("Validated version")
    replacement_payload = _valid_pptx_bytes("Replacement version")
    output = workspace / "brief.pptx"
    output.write_bytes(validated_payload)
    media_root = tmp_path / "media"
    ctx = ToolContext(
        workspace_dir=str(workspace),
        artifact_media_root=str(media_root),
        artifact_session_id="session-1",
        session_key="agent:main:webchat:session-1",
    )
    validate = artifacts_module.validate_artifact_for_delivery

    def validate_then_replace(*args: object, **kwargs: object) -> object:
        report = validate(*args, **kwargs)
        output.write_bytes(replacement_payload)
        return report

    monkeypatch.setattr(
        artifacts_module,
        "validate_artifact_for_delivery",
        validate_then_replace,
    )

    token = current_tool_context.set(ctx)
    try:
        result = json.loads(await publish_artifact(path=output.name, mime=PPTX_MIME))
    finally:
        current_tool_context.reset(token)

    _, material_path = ArtifactStore(media_root).resolve_for_download(
        result["artifact"]["id"],
        session_id="session-1",
    )
    assert output.read_bytes() == replacement_payload
    assert material_path.read_bytes() == validated_payload


@pytest.mark.asyncio
async def test_publish_artifact_tool_hides_local_path_from_non_owner_channel(
    tmp_path: Path,
) -> None:
    workspace = tmp_path / "workspace"
    workspace.mkdir()
    output = workspace / "report.txt"
    output.write_text("ready", encoding="utf-8")
    ctx = ToolContext(
        is_owner=False,
        caller_kind=CallerKind.CHANNEL,
        channel_kind="feishu",
        workspace_dir=str(workspace),
        artifact_media_root=str(tmp_path / "media"),
        artifact_session_id="session-1",
        session_key="agent:main:feishu:direct:u1",
    )

    token = current_tool_context.set(ctx)
    try:
        result = await publish_artifact(path="report.txt", name="final.txt", mime="text/plain")
    finally:
        current_tool_context.reset(token)

    payload = json.loads(result)
    assert payload["status"] == "published"
    assert "download_url" not in payload["artifact"]
    assert "local_path" not in payload["artifact"]
    assert "workspace_path" not in payload["artifact"]
    assert "local_path" not in payload["note"]
    assert "final response" in payload["note"]


@pytest.mark.asyncio
async def test_publish_artifact_tool_accepts_workspace_alias(tmp_path: Path) -> None:
    workspace = tmp_path / "workspace"
    workspace.mkdir()
    output = workspace / "paper.pdf"
    output.write_bytes(b"%PDF-1.5\nready")
    ctx = ToolContext(
        workspace_dir=str(workspace),
        artifact_media_root=str(tmp_path / "media"),
        artifact_session_id="session-1",
        session_key="agent:main:webchat:session-1",
    )

    token = current_tool_context.set(ctx)
    try:
        result = await publish_artifact(
            path="/workspace/paper.pdf",
            name="paper.pdf",
            mime="application/pdf",
        )
    finally:
        current_tool_context.reset(token)

    payload = json.loads(result)
    assert payload["status"] == "published"
    assert payload["artifact"]["name"] == "paper.pdf"
    assert len(ctx.published_artifacts) == 1


@pytest.mark.asyncio
async def test_publish_artifact_tool_is_idempotent_for_existing_turn_artifact(
    tmp_path: Path,
) -> None:
    workspace = tmp_path / "workspace"
    workspace.mkdir()
    output = workspace / "generated-image.png"
    output.write_bytes(b"\x89PNG\r\n\x1a\nsame image")
    ctx = ToolContext(
        workspace_dir=str(workspace),
        artifact_media_root=str(tmp_path / "media"),
        artifact_session_id="session-1",
        session_key="agent:main:feishu:direct:u1",
    )

    token = current_tool_context.set(ctx)
    try:
        first = json.loads(
            await publish_artifact(
                path="generated-image.png",
                name="generated-image.png",
                mime="image/png",
            )
        )
        second = json.loads(
            await publish_artifact(
                path="generated-image.png",
                name="OpenStarry Code-Mascot.png",
                mime="image/png",
            )
        )
    finally:
        current_tool_context.reset(token)

    assert first["status"] == "published"
    assert second["status"] == "already_published"
    assert second["artifact"]["id"] == first["artifact"]["id"]
    assert second["artifact"]["name"] == "generated-image.png"
    assert "already registered" in second["note"]
    assert len(ctx.published_artifacts) == 1


@pytest.mark.asyncio
async def test_publish_artifact_tool_reuses_existing_session_deliverable_across_contexts(
    tmp_path: Path,
) -> None:
    workspace = tmp_path / "workspace"
    workspace.mkdir()
    output = workspace / "brief.pptx"
    output.write_bytes(_valid_pptx_bytes())
    media_root = tmp_path / "media"

    ctx1 = ToolContext(
        workspace_dir=str(workspace),
        artifact_media_root=str(media_root),
        artifact_session_id="session-1",
        session_key="agent:main:webchat:session-1",
    )
    token = current_tool_context.set(ctx1)
    try:
        first = json.loads(
            await publish_artifact(
                path="brief.pptx",
                name="brief.pptx",
                mime=PPTX_MIME,
            )
        )
    finally:
        current_tool_context.reset(token)

    ctx2 = ToolContext(
        workspace_dir=str(workspace),
        artifact_media_root=str(media_root),
        artifact_session_id="session-1",
        session_key="agent:main:webchat:session-1",
    )
    token = current_tool_context.set(ctx2)
    try:
        second = json.loads(
            await publish_artifact(
                path="brief.pptx",
                name="brief.pptx",
                mime=PPTX_MIME,
            )
        )
    finally:
        current_tool_context.reset(token)

    assert first["status"] == "published"
    assert second["status"] == "already_published"
    assert second["artifact"]["id"] == first["artifact"]["id"]
    assert len(ctx1.published_artifacts) == 1
    assert len(ctx2.published_artifacts) == 1
    assert ctx2.published_artifacts[0]["id"] == first["artifact"]["id"]


@pytest.mark.asyncio
async def test_publish_artifact_runs_pptx_read_and_validation_off_the_event_loop(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    import threading

    import openstarry_code.tools.builtin.artifacts as artifacts_module

    workspace = tmp_path / "workspace"
    workspace.mkdir()
    output = workspace / "brief.pptx"
    output.write_bytes(_valid_pptx_bytes("Loop-friendly deck"))
    ctx = ToolContext(
        workspace_dir=str(workspace),
        artifact_media_root=str(tmp_path / "media"),
        artifact_session_id="session-1",
        session_key="agent:main:webchat:session-1",
    )

    validation_threads: list[threading.Thread] = []
    real_validate = artifacts_module.validate_artifact_for_delivery

    def recording_validate(*args: object, **kwargs: object) -> None:
        validation_threads.append(threading.current_thread())
        real_validate(*args, **kwargs)  # type: ignore[arg-type]

    monkeypatch.setattr(
        artifacts_module,
        "validate_artifact_for_delivery",
        recording_validate,
    )
    loop_thread = threading.current_thread()
    token = current_tool_context.set(ctx)
    try:
        result = json.loads(await publish_artifact(path=output.name, mime=PPTX_MIME))
    finally:
        current_tool_context.reset(token)

    assert result["status"] == "published"
    assert validation_threads
    assert all(thread is not loop_thread for thread in validation_threads)


@pytest.mark.asyncio
async def test_publish_artifact_tool_republishes_changed_bytes_at_same_path(
    tmp_path: Path,
) -> None:
    workspace = tmp_path / "workspace"
    workspace.mkdir()
    output = workspace / "report.txt"
    output.write_text("first", encoding="utf-8")
    ctx = ToolContext(
        workspace_dir=str(workspace),
        artifact_media_root=str(tmp_path / "media"),
        artifact_session_id="session-1",
        session_key="agent:main:webchat:session-1",
    )

    token = current_tool_context.set(ctx)
    try:
        first = json.loads(await publish_artifact(path="report.txt", mime="text/plain"))
        output.write_text("second", encoding="utf-8")
        second = json.loads(await publish_artifact(path="report.txt", mime="text/plain"))
    finally:
        current_tool_context.reset(token)

    assert first["status"] == "published"
    assert second["status"] == "published"
    assert second["artifact"]["id"] != first["artifact"]["id"]
    assert len(ctx.published_artifacts) == 2


@pytest.mark.asyncio
async def test_publish_artifact_tool_reports_storage_write_failure(
    monkeypatch, tmp_path: Path
) -> None:
    workspace = tmp_path / "workspace"
    workspace.mkdir()
    output = workspace / "report.txt"
    output.write_text("ready", encoding="utf-8")
    ctx = ToolContext(
        workspace_dir=str(workspace),
        artifact_media_root=str(tmp_path / "media"),
        artifact_session_id="session-1",
        session_key="agent:main:webchat:session-1",
    )

    def fail_publish_file(*args: object, **kwargs: object) -> None:
        raise FileNotFoundError("media temp path unavailable")

    monkeypatch.setattr(ArtifactStore, "publish_file", fail_publish_file)
    token = current_tool_context.set(ctx)
    try:
        with pytest.raises(ToolError, match="artifact storage path is unavailable"):
            await publish_artifact(path="report.txt", name="final.txt", mime="text/plain")
    finally:
        current_tool_context.reset(token)


@pytest.mark.asyncio
async def test_publish_artifact_tool_missing_file_reports_workspace_candidates(
    tmp_path: Path,
) -> None:
    workspace = tmp_path / "workspace"
    reports = workspace / "reports"
    reports.mkdir(parents=True)
    candidate = reports / "AI Agent Comparison 2026.pptx"
    candidate.write_bytes(b"pptx")
    ctx = ToolContext(
        workspace_dir=str(workspace),
        artifact_media_root=str(tmp_path / "media"),
        artifact_session_id="session-1",
        session_key="agent:main:webchat:session-1",
    )

    token = current_tool_context.set(ctx)
    try:
        with pytest.raises(ToolError) as exc_info:
            await publish_artifact(path="AI_Agent_Comparison_2026.pptx")
    finally:
        current_tool_context.reset(token)

    message = str(exc_info.value)
    assert "artifact file not found" in message
    assert f"active workspace: {workspace.resolve()}" in message
    assert "resolved path:" in message
    assert "candidate files:" in message
    assert "reports/AI Agent Comparison 2026.pptx" in message.replace("\\", "/")


@pytest.mark.asyncio
async def test_publish_artifact_rejects_foreign_posix_target_with_workspace_hint(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    import openstarry_code.tools.builtin.artifacts as artifacts_module

    monkeypatch.setattr(artifacts_module, "os", SimpleNamespace(name="nt"), raising=False)
    workspace = tmp_path / "workspace"
    workspace.mkdir()
    actual = workspace / "report.pptx"
    actual.write_bytes(b"pptx")
    ctx = ToolContext(
        workspace_dir=str(workspace),
        artifact_media_root=str(tmp_path / "media"),
        artifact_session_id="session-1",
        session_key="agent:main:webchat:session-1",
    )

    token = current_tool_context.set(ctx)
    try:
        with pytest.raises(ToolError) as exc_info:
            await publish_artifact(path="/Users/a1/Desktop/report.pptx")
    finally:
        current_tool_context.reset(token)

    message = str(exc_info.value)
    assert "foreign_host_path" in message
    assert "requested path is from another host/platform" in message
    assert "report.pptx" in message
    assert "D:\\Users" not in message
    assert not ctx.published_artifacts


@pytest.mark.asyncio
async def test_publish_artifact_tool_rejects_missing_workspace_and_escape(tmp_path: Path) -> None:
    workspace = tmp_path / "workspace"
    workspace.mkdir()
    outside = tmp_path / "outside.txt"
    outside.write_text("no", encoding="utf-8")

    token = current_tool_context.set(
        ToolContext(
            artifact_media_root=str(tmp_path / "media"),
            artifact_session_id="session-1",
            session_key="agent:main:webchat:session-1",
        )
    )
    try:
        with pytest.raises(ToolError):
            await publish_artifact(path=str(outside))
    finally:
        current_tool_context.reset(token)

    token = current_tool_context.set(
        ToolContext(
            workspace_dir=str(workspace),
            artifact_media_root=str(tmp_path / "media"),
            artifact_session_id="session-1",
            session_key="agent:main:webchat:session-1",
        )
    )
    try:
        with pytest.raises(ToolError):
            await publish_artifact(path="../outside.txt")
    finally:
        current_tool_context.reset(token)


def test_copy_session_artifacts_rebinds_to_child_and_preserves_isolation(
    tmp_path: Path,
) -> None:
    store = ArtifactStore(tmp_path)
    ref = store.publish_bytes(
        b"deliverable bytes",
        session_id="parent-1",
        session_key="agent:main:webchat:parent-1",
        name="report.txt",
        mime="text/plain",
        source="publish_artifact",
    )

    # Before the copy the child cannot see the parent's artifact.
    with pytest.raises(ArtifactNotFoundError):
        store.resolve_for_download(ref.id, session_id="child-1")

    copied = store.copy_session_artifacts(
        source_session_id="parent-1",
        target_session_id="child-1",
        target_session_key="agent:main:webchat:child-1",
    )
    assert copied == 1

    child_ref, child_path = store.resolve_for_download(ref.id, session_id="child-1")
    assert child_path.read_bytes() == b"deliverable bytes"
    assert child_ref.id == ref.id  # stable id keeps the transcript/URL linkage valid
    assert child_ref.session_id == "child-1"
    assert child_ref.session_key == "agent:main:webchat:child-1"
    assert child_ref.sha256 == ref.sha256
    child_page = store.list_refs(session_id="child-1", limit=10)
    assert [listed.id for listed in child_page.refs] == [ref.id]

    # The parent still owns its copy and an unrelated session stays blocked.
    parent_ref, _ = store.resolve_for_download(ref.id, session_id="parent-1")
    assert parent_ref.session_id == "parent-1"
    with pytest.raises(ArtifactNotFoundError):
        store.resolve_for_download(ref.id, session_id="stranger")

    # Re-copying is idempotent: nothing new is materialized.
    assert (
        store.copy_session_artifacts(
            source_session_id="parent-1",
            target_session_id="child-1",
            target_session_key="agent:main:webchat:child-1",
        )
        == 0
    )


def test_copy_session_artifacts_carries_thumbnail(tmp_path: Path) -> None:
    from PIL import Image

    store = ArtifactStore(tmp_path)
    out = io.BytesIO()
    Image.new("RGB", (8, 8), color="red").save(out, format="PNG")
    ref = store.publish_bytes(
        out.getvalue(),
        session_id="parent-1",
        session_key="agent:main:webchat:parent-1",
        name="chart.png",
        mime="image/png",
        source="publish_artifact",
    )
    assert ref.has_thumbnail is True

    store.copy_session_artifacts(
        source_session_id="parent-1",
        target_session_id="child-1",
        target_session_key="agent:main:webchat:child-1",
    )

    thumbnail = store.resolve_thumbnail_for_download(ref.id, session_id="child-1")
    assert thumbnail is not None
    _, thumb_path = thumbnail
    assert thumb_path.exists()


def test_copy_session_artifacts_reads_legacy_short_layout(tmp_path: Path) -> None:
    store = ArtifactStore(tmp_path)
    session_id = "532d5065-abce-499f-97b0-bbf2a067d5ab"
    ref = store.publish_bytes(
        b"legacy material",
        session_id=session_id,
        session_key="agent:main:webchat:legacy",
        name="old.txt",
        mime="text/plain",
        source="publish_artifact",
    )

    # Relocate the artifact into the 16-char legacy session/artifact layout.
    current_dir = store.path_for(ref).parent
    legacy_session_token = hashlib.sha256(session_id.encode("utf-8")).hexdigest()[:16]
    legacy_artifact_token = hashlib.sha256(ref.id.encode("utf-8")).hexdigest()[:16]
    legacy_dir = tmp_path / "artifacts" / "s" / legacy_session_token / legacy_artifact_token
    legacy_dir.parent.mkdir(parents=True)
    current_dir.rename(legacy_dir)

    copied = store.copy_session_artifacts(
        source_session_id=session_id,
        target_session_id="child-1",
        target_session_key="agent:main:webchat:child-1",
    )
    assert copied == 1
    _, child_path = store.resolve_for_download(ref.id, session_id="child-1")
    assert child_path.read_bytes() == b"legacy material"


def test_copy_session_artifacts_reads_legacy_plain_layout(tmp_path: Path) -> None:
    from openstarry_code.artifacts import _safe_token

    store = ArtifactStore(tmp_path)
    session_id = "plain-session"
    ref = store.publish_bytes(
        b"plain layout material",
        session_id=session_id,
        session_key="agent:main:webchat:plain",
        name="legacy.txt",
        mime="text/plain",
        source="publish_artifact",
    )

    # Relocate into the oldest "plain" layout where the material file is named by the
    # sha (not "data"): artifacts/<safe_token(session)>/<artifact-id>/<sha256>.
    current_dir = store.path_for(ref).parent
    plain_dir = tmp_path / "artifacts" / _safe_token(session_id) / ref.id
    plain_dir.mkdir(parents=True)
    (current_dir / "data").rename(plain_dir / ref.sha256)
    (current_dir / "meta.json").rename(plain_dir / "meta.json")

    copied = store.copy_session_artifacts(
        source_session_id=session_id,
        target_session_id="child-1",
        target_session_key="agent:main:webchat:child-1",
    )
    assert copied == 1
    _, child_path = store.resolve_for_download(ref.id, session_id="child-1")
    assert child_path.read_bytes() == b"plain layout material"


def test_copy_session_artifacts_skips_artifact_with_missing_material(tmp_path: Path) -> None:
    store = ArtifactStore(tmp_path)
    good = store.publish_bytes(
        b"good bytes",
        session_id="parent-1",
        session_key="agent:main:webchat:parent-1",
        name="good.txt",
        mime="text/plain",
        source="publish_artifact",
    )
    bad = store.publish_bytes(
        b"vanishing bytes",
        session_id="parent-1",
        session_key="agent:main:webchat:parent-1",
        name="bad.txt",
        mime="text/plain",
        source="publish_artifact",
    )
    # Drop the bad artifact's material, leaving its meta.json behind.
    store.path_for(bad).unlink()

    copied = store.copy_session_artifacts(
        source_session_id="parent-1",
        target_session_id="child-1",
        target_session_key="agent:main:webchat:child-1",
    )
    assert copied == 1  # only the artifact with intact material is carried
    _, child_path = store.resolve_for_download(good.id, session_id="child-1")
    assert child_path.read_bytes() == b"good bytes"
    with pytest.raises(ArtifactNotFoundError):
        store.resolve_for_download(bad.id, session_id="child-1")


def test_strip_artifact_markers_preserves_surrounding_whitespace() -> None:
    marker = "[generated artifact omitted: report.html (text/html)]"

    assert strip_artifact_markers_from_text(f"line1\n{marker}\nline2") in (
        "line1\nline2",
        "line1\n\nline2",
    )
    assert (
        strip_artifact_markers_from_text(f"Here is the summary. {marker} Let me know.")
        == "Here is the summary. Let me know."
    )


def test_strip_artifact_markers_handles_bracket_in_name() -> None:
    marker = artifact_marker({"name": "weird].html", "mime": "text/html"})
    assert marker == "[generated artifact omitted: weird].html (text/html)]"

    cleaned = strip_artifact_markers_from_text(f"Done!\n{marker}\nAnything else?")
    assert "]" not in cleaned
    assert ".html" not in cleaned
